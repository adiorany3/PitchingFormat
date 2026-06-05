import html
import re
from io import BytesIO
from typing import Any

import streamlit as st

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt

from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import cm
from reportlab.platypus import (
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table as RLTable,
    TableStyle,
)


st.set_page_config(
    page_title="Seed Investor Pitch Deck Generator",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded",
)


# ==============================
# Streamlit Chrome Removal
# ==============================
def hide_streamlit_emblems():
    """Hide Streamlit chrome and apply theme-aware readable UI styling."""
    st.markdown(
        """
        <style>
            :root {
                --deck-bg: var(--background-color, #f8fafc);
                --deck-surface: var(--secondary-background-color, #ffffff);
                --deck-surface-soft: rgba(148, 163, 184, 0.10);
                --deck-surface-strong: rgba(148, 163, 184, 0.16);
                --deck-text: var(--text-color, #0f172a);
                --deck-primary: var(--primary-color, #2563eb);
                --deck-muted: rgba(100, 116, 139, 0.96);
                --deck-muted: color-mix(in srgb, var(--text-color, #0f172a) 68%, transparent);
                --deck-border: rgba(148, 163, 184, 0.35);
                --deck-border-strong: rgba(148, 163, 184, 0.55);
                --deck-shadow: 0 14px 35px rgba(15, 23, 42, 0.08);
                --deck-radius: 16px;
            }

            /* Remove Streamlit chrome / emblem */
            #MainMenu {visibility: hidden !important;}
            footer {visibility: hidden !important;}
            header {visibility: hidden !important;}
            [data-testid="stHeader"] {display: none !important;}
            [data-testid="stToolbar"] {display: none !important;}
            [data-testid="stDecoration"] {display: none !important;}
            [data-testid="stStatusWidget"] {display: none !important;}
            [data-testid="stDeployButton"] {display: none !important;}
            [data-testid="collapsedControl"] {display: none !important;}
            .stAppDeployButton {display: none !important;}
            .st-emotion-cache-1dp5vir {display: none !important;}
            .st-emotion-cache-14xtw13 {display: none !important;}
            .viewerBadge_container__1QSob,
            .viewerBadge_link__1S137,
            .viewerBadge_text__1JaDK {display: none !important;}

            /* App layout */
            .stApp {
                background: var(--deck-bg) !important;
                color: var(--deck-text) !important;
            }
            .block-container {
                max-width: 1320px;
                padding-top: 1.35rem !important;
                padding-bottom: 5.8rem !important;
            }
            [data-testid="stSidebar"] > div:first-child {
                background: var(--deck-surface) !important;
                border-right: 1px solid var(--deck-border);
            }
            [data-testid="stSidebar"] * {
                color: var(--deck-text);
            }

            /* Text readability across light/dark themes */
            h1, h2, h3, h4, h5, h6,
            [data-testid="stMarkdownContainer"] p,
            [data-testid="stMarkdownContainer"] li,
            [data-testid="stCaptionContainer"],
            label {
                color: var(--deck-text) !important;
            }
            small, .caption, [data-testid="stCaptionContainer"] {
                color: var(--deck-muted) !important;
            }
            hr {
                border-color: var(--deck-border) !important;
            }

            /* Tabs */
            [data-testid="stTabs"] [role="tablist"] {
                gap: 0.35rem;
                border-bottom: 1px solid var(--deck-border);
                flex-wrap: wrap;
            }
            [data-testid="stTabs"] button[role="tab"] {
                min-height: 42px;
                padding: 0.55rem 0.85rem;
                border-radius: 999px 999px 0 0;
                color: var(--deck-muted) !important;
                background: transparent !important;
                border: 1px solid transparent !important;
            }
            [data-testid="stTabs"] button[role="tab"][aria-selected="true"] {
                color: var(--deck-text) !important;
                background: var(--deck-surface) !important;
                border: 1px solid var(--deck-border) !important;
                border-bottom-color: var(--deck-surface) !important;
                font-weight: 700;
            }

            /* Form controls */
            [data-testid="stTextInput"] input,
            [data-testid="stNumberInput"] input,
            [data-testid="stTextArea"] textarea,
            [data-baseweb="select"] > div,
            [data-testid="stFileUploader"] section,
            [data-testid="stColorPicker"] input {
                background: var(--deck-surface) !important;
                color: var(--deck-text) !important;
                border-color: var(--deck-border) !important;
                caret-color: var(--deck-primary) !important;
            }
            [data-testid="stTextInput"] input::placeholder,
            [data-testid="stTextArea"] textarea::placeholder {
                color: var(--deck-muted) !important;
                opacity: 0.82;
            }
            [data-testid="stTextInput"] input:focus,
            [data-testid="stNumberInput"] input:focus,
            [data-testid="stTextArea"] textarea:focus,
            [data-baseweb="select"] > div:focus-within {
                border-color: var(--deck-primary) !important;
                box-shadow: 0 0 0 2px color-mix(in srgb, var(--deck-primary) 24%, transparent) !important;
            }
            [data-testid="stWidgetLabel"] p {
                color: var(--deck-text) !important;
                font-weight: 650;
            }
            [data-testid="InputInstructions"],
            [data-testid="stWidgetLabel"] small,
            [data-testid="stHelp"] {
                color: var(--deck-muted) !important;
            }

            /* Buttons */
            .stButton > button,
            .stDownloadButton > button {
                border-radius: 12px !important;
                border: 1px solid var(--deck-border-strong) !important;
                background: var(--deck-surface) !important;
                color: var(--deck-text) !important;
                font-weight: 700 !important;
                min-height: 44px;
            }
            .stButton > button[kind="primary"],
            .stButton > button[data-testid="baseButton-primary"] {
                background: var(--deck-primary) !important;
                color: #ffffff !important;
                border-color: var(--deck-primary) !important;
            }
            .stDownloadButton > button {
                background: var(--deck-primary) !important;
                color: #ffffff !important;
                border-color: var(--deck-primary) !important;
            }

            /* Native Streamlit metric readability */
            [data-testid="stMetric"] {
                background: var(--deck-surface) !important;
                border: 1px solid var(--deck-border);
                border-radius: var(--deck-radius);
                padding: 0.9rem 1rem;
                box-shadow: var(--deck-shadow);
            }
            [data-testid="stMetricLabel"] p {
                color: var(--deck-muted) !important;
                font-weight: 700;
            }
            [data-testid="stMetricValue"] {
                color: var(--deck-text) !important;
            }

            /* Custom reusable blocks */
            .developer-footer {
                position: fixed;
                left: 0;
                bottom: 0;
                width: 100%;
                background: color-mix(in srgb, var(--deck-surface) 94%, transparent);
                backdrop-filter: blur(10px);
                border-top: 1px solid var(--deck-border);
                color: var(--deck-muted) !important;
                text-align: center;
                padding: 8px 0;
                font-size: 12px;
                line-height: 1.4;
                z-index: 9999;
            }
            .guide-box,
            .insight-card,
            .readable-panel {
                background: var(--deck-surface);
                border: 1px solid var(--deck-border);
                border-radius: var(--deck-radius);
                padding: 15px 17px;
                margin: 8px 0 18px 0;
                color: var(--deck-text) !important;
                box-shadow: var(--deck-shadow);
            }
            .guide-box {
                border-left: 4px solid var(--deck-primary);
            }
            .guide-box strong,
            .guide-box .guide-title,
            .insight-card h4 {
                color: var(--deck-text) !important;
                font-weight: 800;
            }
            .guide-box p,
            .guide-box span,
            .insight-card p,
            .insight-card li,
            .readable-panel p,
            .readable-panel li {
                color: var(--deck-muted) !important;
                font-size: 0.94rem;
                line-height: 1.58;
            }
            .insight-card h4 {
                margin: 0 0 10px 0;
                font-size: 1.02rem;
            }
            .insight-card ul {
                margin: 0.35rem 0 0 1.1rem;
                padding: 0;
            }
            .insight-card li {
                margin-bottom: 0.48rem;
            }

            /* Alerts and upload blocks should keep contrast */
            [data-testid="stAlert"] {
                border-radius: var(--deck-radius) !important;
                border: 1px solid var(--deck-border) !important;
            }
            [data-testid="stFileUploader"] section {
                border-radius: var(--deck-radius) !important;
            }

            @media (max-width: 900px) {
                .block-container {
                    padding-left: 1rem !important;
                    padding-right: 1rem !important;
                }
                .developer-footer {
                    font-size: 11px;
                    padding: 7px 0;
                }
                [data-testid="stTabs"] button[role="tab"] {
                    min-height: 38px;
                    padding: 0.45rem 0.65rem;
                }
            }
        </style>
        """,
        unsafe_allow_html=True,
    )

hide_streamlit_emblems()


# ==============================
# Professional Seed Deck Theme
# ==============================
THEME = {
    "bg": RGBColor(248, 250, 252),
    "white": RGBColor(255, 255, 255),
    "ink": RGBColor(15, 23, 42),
    "muted": RGBColor(100, 116, 139),
    "line": RGBColor(226, 232, 240),
    "dark": RGBColor(2, 6, 23),
    "green": RGBColor(22, 163, 74),
    "amber": RGBColor(217, 119, 6),
    "red": RGBColor(220, 38, 38),
    "blue_soft": RGBColor(239, 246, 255),
    "blue_line": RGBColor(191, 219, 254),
}

FONT_HEAD = "Aptos Display"
FONT_BODY = "Aptos"
DEVELOPER_FOOTER = "Developed by Galuh Adi Insani"

SLIDE_W = 13.333
SLIDE_H = 7.5


# ==============================
# Utility Helpers
# ==============================
def rgb(hex_color: str) -> RGBColor:
    value = (hex_color or "#2563EB").replace("#", "").strip()

    if len(value) != 6:
        value = "2563EB"

    return RGBColor(
        int(value[:2], 16),
        int(value[2:4], 16),
        int(value[4:], 16),
    )


def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def lines(text: str, limit: int = 5) -> list[str]:
    result = []

    for row in (text or "").splitlines():
        cleaned = row.strip().lstrip("-•*0123456789. ").strip()

        if cleaned:
            result.append(cleaned)

    return result[:limit]


def money(value: int | float, currency: str) -> str:
    if currency == "Rp":
        return f"Rp {value:,.0f}".replace(",", ".")

    return f"{currency} {value:,.0f}"


def filename(value: str) -> str:
    value = re.sub(
        r"[^a-zA-Z0-9 _-]",
        "",
        value or "pitch-deck",
    )

    return value.strip().lower().replace(" ", "-") or "pitch-deck"


def parse_percent(value: str) -> float | None:
    if not value:
        return None

    match = re.search(r"-?\d+(?:[\.,]\d+)?", str(value))

    if not match:
        return None

    return float(match.group(0).replace(",", "."))


def safe_div(numerator: float, denominator: float) -> float | None:
    if denominator == 0:
        return None

    return numerator / denominator


def pct(value: float | None) -> str:
    if value is None:
        return "N/A"

    return f"{value * 100:.1f}%"


def multiple(value: float | None) -> str:
    if value is None:
        return "N/A"

    return f"{value:.1f}x"


def short_money(value: float, currency: str) -> str:
    sign = "-" if value < 0 else ""
    value = abs(value)

    if currency == "Rp":
        if value >= 1_000_000_000_000:
            return f"{sign}Rp {value / 1_000_000_000_000:.1f}T"
        if value >= 1_000_000_000:
            return f"{sign}Rp {value / 1_000_000_000:.1f}M"
        if value >= 1_000_000:
            return f"{sign}Rp {value / 1_000_000:.1f}Jt"
        return f"{sign}Rp {value:,.0f}".replace(",", ".")

    if value >= 1_000_000_000:
        return f"{sign}{currency} {value / 1_000_000_000:.1f}B"
    if value >= 1_000_000:
        return f"{sign}{currency} {value / 1_000_000:.1f}M"
    if value >= 1_000:
        return f"{sign}{currency} {value / 1_000:.1f}K"

    return f"{sign}{currency} {value:,.0f}"



# ==============================
# Startup Pitch Glossary
# ==============================
STARTUP_TERMS = [
    {
        "category": "Fundraising",
        "term": "Pre-Seed Round",
        "simple": "Pendanaan paling awal untuk membuktikan problem, membangun MVP, dan mencari validasi pasar pertama.",
        "formula": "Tidak ada rumus baku. Biasanya ditentukan dari kebutuhan dana 6-18 bulan + milestone validasi awal.",
        "example": "Butuh Rp 600 juta untuk 12 bulan membangun MVP, pilot 20 customer, dan validasi pricing.",
        "used_in": "Identitas, Fundraising Ask, Milestones",
    },
    {
        "category": "Fundraising",
        "term": "Seed Round",
        "simple": "Pendanaan untuk mempercepat produk, akuisisi pelanggan, traction awal, dan pembuktian model bisnis.",
        "formula": "Ask seed = kebutuhan runway 12-18 bulan + biaya mencapai milestone berikutnya.",
        "example": "Ask Rp 1,5M untuk 18 bulan sampai Rp 500 juta MRR dan 8.000 active businesses.",
        "used_in": "Cover, Financials, Fundraising Ask",
    },
    {
        "category": "Fundraising",
        "term": "Ask / Jumlah Pendanaan",
        "simple": "Jumlah dana yang diminta dari investor pada round ini.",
        "formula": "Ask = estimasi biaya bulanan x target runway + buffer eksekusi.",
        "example": "Burn Rp 80 juta/bulan x 18 bulan = Rp 1,44M. Ask dapat dibulatkan menjadi Rp 1,5M.",
        "used_in": "Cover, Financials, Fundraising Ask",
    },
    {
        "category": "Fundraising",
        "term": "Use of Funds",
        "simple": "Rencana penggunaan dana investor dalam kategori utama.",
        "formula": "Total use of funds harus = 100% dari ask. Contoh kategori: Product, Sales & Marketing, Team, Operations.",
        "example": "Product 40%, Sales & Marketing 35%, Customer Success 15%, Operations 10%.",
        "used_in": "Fundraising Ask",
    },
    {
        "category": "Fundraising",
        "term": "Runway",
        "simple": "Berapa bulan startup bisa berjalan sebelum dana habis.",
        "formula": "Runway = cash tersedia / net burn bulanan.",
        "example": "Cash Rp 1,5M dan burn Rp 83 juta/bulan -> runway sekitar 18 bulan.",
        "used_in": "Financials, Fundraising Ask",
    },
    {
        "category": "Fundraising",
        "term": "Milestone",
        "simple": "Target penting yang harus dicapai agar startup naik kelas atau siap fundraising berikutnya.",
        "formula": "Milestone baik harus punya periode + target + success metric + owner.",
        "example": "0-6 bulan: rilis MVP v2, success metric: 500 paid users, owner: Product + Growth.",
        "used_in": "Milestones, Financials, Fundraising Ask",
    },
    {
        "category": "Fundraising",
        "term": "Success Metric",
        "simple": "Ukuran keberhasilan yang membuat milestone tidak subjektif.",
        "formula": "Gunakan angka: MRR, ARR, active user, retention, CAC payback, gross margin, partnership, atau launch date.",
        "example": "Target bukan 'marketing lebih kuat', tetapi 'CAC payback < 3 bulan dari channel komunitas'.",
        "used_in": "Milestones",
    },
    {
        "category": "Fundraising",
        "term": "Valuation",
        "simple": "Nilai perusahaan yang menjadi dasar investor membeli saham.",
        "formula": "Post-money valuation = pre-money valuation + dana masuk. Dilution = dana masuk / post-money valuation.",
        "example": "Pre-money Rp 10M + dana Rp 2M = post-money Rp 12M. Dilution investor = 2/12 = 16,7%.",
        "used_in": "Q&A investor, tidak selalu masuk slide seed awal",
    },
    {
        "category": "Fundraising",
        "term": "Dilution",
        "simple": "Persentase kepemilikan founder yang berkurang karena saham baru diterbitkan untuk investor.",
        "formula": "Dilution investor = dana masuk / post-money valuation.",
        "example": "Investor masuk Rp 1,5M dengan post-money Rp 10M -> dilution 15%.",
        "used_in": "Q&A investor",
    },
    {
        "category": "Market",
        "term": "TAM - Total Addressable Market",
        "simple": "Total peluang pasar jika startup bisa melayani seluruh pasar yang relevan.",
        "formula": "TAM = jumlah seluruh target customer x potensi revenue per customer per tahun.",
        "example": "1 juta UMKM x Rp 1,2 juta/tahun = Rp 1,2T TAM.",
        "used_in": "Market",
    },
    {
        "category": "Market",
        "term": "SAM - Serviceable Available Market",
        "simple": "Bagian TAM yang realistis dijangkau oleh produk, wilayah, channel, dan model bisnis saat ini.",
        "formula": "SAM = jumlah customer dalam segmen yang bisa dilayani x revenue per customer per tahun.",
        "example": "150 ribu UMKM F&B yang sudah digital x Rp 1,2 juta/tahun = Rp 180M SAM.",
        "used_in": "Market",
    },
    {
        "category": "Market",
        "term": "SOM - Serviceable Obtainable Market",
        "simple": "Bagian SAM yang realistis dimenangkan dalam 2-5 tahun pertama.",
        "formula": "SOM = target market share realistis x SAM.",
        "example": "Target 5% dari SAM Rp 180M = Rp 9M SOM awal.",
        "used_in": "Market",
    },
    {
        "category": "Market",
        "term": "ICP - Ideal Customer Profile",
        "simple": "Profil customer yang paling cocok, paling butuh solusi, dan paling mudah dijangkau.",
        "formula": "ICP = segmen + ukuran bisnis + masalah utama + kemampuan bayar + channel akses.",
        "example": "UMKM F&B dengan 3-50 transaksi/hari, masih memakai spreadsheet, aktif di komunitas WhatsApp.",
        "used_in": "Go-To-Market",
    },
    {
        "category": "Market",
        "term": "Wedge Market",
        "simple": "Segmen awal yang kecil tetapi tajam untuk dimenangkan sebelum ekspansi ke pasar lebih besar.",
        "formula": "Pilih segmen dengan pain tinggi, akses mudah, sales cycle pendek, dan use case jelas.",
        "example": "Mulai dari UMKM F&B, lalu ekspansi ke ritel, jasa, dan wholesale.",
        "used_in": "Market, GTM",
    },
    {
        "category": "Revenue",
        "term": "Revenue",
        "simple": "Pendapatan yang benar-benar diterima atau diakui dari customer.",
        "formula": "Revenue = jumlah customer berbayar x harga rata-rata dalam periode tertentu.",
        "example": "500 customer x Rp 100.000/bulan = Rp 50 juta revenue bulanan.",
        "used_in": "Traction, Financials",
    },
    {
        "category": "Revenue",
        "term": "MRR - Monthly Recurring Revenue",
        "simple": "Pendapatan berulang per bulan, biasanya untuk subscription/SaaS.",
        "formula": "MRR = jumlah pelanggan aktif berbayar x biaya subscription bulanan rata-rata.",
        "example": "850 customer x Rp 100.000/bulan = Rp 85 juta MRR.",
        "used_in": "Traction, Milestones, Financials",
    },
    {
        "category": "Revenue",
        "term": "ARR - Annual Recurring Revenue",
        "simple": "Pendapatan berulang tahunan dari model subscription.",
        "formula": "ARR = MRR x 12.",
        "example": "MRR Rp 85 juta x 12 = Rp 1,02M ARR.",
        "used_in": "Traction, Financials",
    },
    {
        "category": "Revenue",
        "term": "GMV - Gross Merchandise Value",
        "simple": "Total nilai transaksi yang lewat platform, belum tentu menjadi revenue perusahaan.",
        "formula": "GMV = jumlah transaksi x nilai transaksi rata-rata.",
        "example": "10.000 transaksi x Rp 150.000 = Rp 1,5M GMV. Jika take rate 5%, revenue = Rp 75 juta.",
        "used_in": "Traction, Marketplace model",
    },
    {
        "category": "Revenue",
        "term": "Take Rate",
        "simple": "Persentase GMV yang menjadi revenue startup.",
        "formula": "Take rate = revenue / GMV.",
        "example": "Revenue Rp 75 juta dari GMV Rp 1,5M -> take rate 5%.",
        "used_in": "Business Model",
    },
    {
        "category": "Unit Economics",
        "term": "ARPU - Average Revenue Per User",
        "simple": "Pendapatan rata-rata dari satu user/customer dalam periode tertentu.",
        "formula": "ARPU = total revenue / jumlah user atau customer berbayar.",
        "example": "Revenue Rp 50 juta/bulan dari 500 customer -> ARPU Rp 100.000/bulan.",
        "used_in": "Business Model, Market, Financials",
    },
    {
        "category": "Unit Economics",
        "term": "Gross Margin",
        "simple": "Persentase revenue yang tersisa setelah biaya langsung melayani customer.",
        "formula": "Gross margin = (revenue - COGS) / revenue x 100%.",
        "example": "Revenue Rp 100 juta, COGS Rp 22 juta -> gross margin 78%.",
        "used_in": "Business Model, Financials",
    },
    {
        "category": "Unit Economics",
        "term": "COGS - Cost of Goods Sold",
        "simple": "Biaya langsung untuk menyediakan produk/jasa kepada customer.",
        "formula": "COGS = biaya server, payment fee, support langsung, bahan baku, atau biaya layanan yang naik mengikuti customer.",
        "example": "Revenue Rp 100 juta, biaya server dan payment fee Rp 10 juta, support langsung Rp 12 juta -> COGS Rp 22 juta.",
        "used_in": "Business Model, Financials",
    },
    {
        "category": "Unit Economics",
        "term": "CAC - Customer Acquisition Cost",
        "simple": "Biaya rata-rata untuk mendapatkan satu customer baru.",
        "formula": "CAC = total biaya sales & marketing / jumlah customer baru.",
        "example": "Biaya marketing Rp 14 juta menghasilkan 100 customer baru -> CAC Rp 140.000.",
        "used_in": "Business Model, GTM, Q&A investor",
    },
    {
        "category": "Unit Economics",
        "term": "CAC Payback",
        "simple": "Waktu yang dibutuhkan untuk mengembalikan biaya akuisisi customer dari gross profit customer tersebut.",
        "formula": "CAC payback = CAC / (ARPU x gross margin bulanan).",
        "example": "CAC Rp 140.000, ARPU Rp 100.000, gross margin 70% -> payback 2 bulan.",
        "used_in": "Business Model, GTM",
    },
    {
        "category": "Unit Economics",
        "term": "LTV - Lifetime Value",
        "simple": "Estimasi gross profit dari satu customer selama customer masih memakai produk.",
        "formula": "LTV sederhana = ARPU bulanan x gross margin x umur customer dalam bulan.",
        "example": "ARPU Rp 100.000, margin 70%, lifetime 12 bulan -> LTV Rp 840.000.",
        "used_in": "Q&A investor, Unit economics",
    },
    {
        "category": "Unit Economics",
        "term": "LTV/CAC",
        "simple": "Rasio nilai customer dibanding biaya mendapatkannya.",
        "formula": "LTV/CAC = LTV / CAC.",
        "example": "LTV Rp 840.000 dan CAC Rp 140.000 -> LTV/CAC = 6x.",
        "used_in": "Business Model, Q&A investor",
    },
    {
        "category": "Traction",
        "term": "Traction",
        "simple": "Bukti bahwa market mulai menerima produk.",
        "formula": "Tidak satu rumus. Gunakan bukti: revenue, paid users, active usage, retention, pilot, LOI, partnership, pipeline.",
        "example": "Rp 85 juta MRR, 1.200 active users, retention D30 72%, growth 28% MoM.",
        "used_in": "Traction",
    },
    {
        "category": "Traction",
        "term": "Active Users",
        "simple": "User yang benar-benar aktif memakai produk dalam periode tertentu.",
        "formula": "DAU = daily active users, WAU = weekly active users, MAU = monthly active users.",
        "example": "2.000 MAU berarti 2.000 user aktif minimal sekali dalam 30 hari terakhir.",
        "used_in": "Traction",
    },
    {
        "category": "Traction",
        "term": "Retention D7 / D30",
        "simple": "Persentase user yang kembali memakai produk setelah 7 atau 30 hari.",
        "formula": "Retention = user yang kembali pada hari ke-N / user yang mulai pada hari pertama x 100%.",
        "example": "100 user daftar, 72 masih aktif hari ke-30 -> D30 retention 72%.",
        "used_in": "Traction, Q&A investor",
    },
    {
        "category": "Traction",
        "term": "Churn",
        "simple": "Persentase customer yang berhenti memakai atau berhenti membayar.",
        "formula": "Customer churn = customer hilang dalam periode / customer awal periode x 100%.",
        "example": "Awal bulan 500 customer, 25 berhenti -> churn 5% per bulan.",
        "used_in": "Traction, Q&A investor",
    },
    {
        "category": "Traction",
        "term": "MoM Growth",
        "simple": "Pertumbuhan dari bulan ke bulan.",
        "formula": "MoM growth = (nilai bulan ini - nilai bulan lalu) / nilai bulan lalu x 100%.",
        "example": "MRR naik dari Rp 70 juta ke Rp 85 juta -> growth 21,4% MoM.",
        "used_in": "Traction, Financials",
    },
    {
        "category": "Traction",
        "term": "Pipeline",
        "simple": "Calon customer atau potensi deal yang sedang diproses sales.",
        "formula": "Pipeline value = jumlah prospek x estimasi nilai deal x probabilitas closing.",
        "example": "20 prospek x Rp 10 juta x peluang 40% = expected pipeline Rp 80 juta.",
        "used_in": "Traction, GTM",
    },
    {
        "category": "Traction",
        "term": "LOI - Letter of Intent",
        "simple": "Surat minat dari calon customer/partner yang menunjukkan niat bekerja sama atau membeli.",
        "formula": "Tidak ada rumus. Nilainya lebih kuat jika ada jumlah, timeline, dan pihak yang jelas.",
        "example": "5 LOI dari koperasi UMKM dengan potensi 2.000 merchant.",
        "used_in": "Traction",
    },
    {
        "category": "GTM",
        "term": "GTM - Go-To-Market",
        "simple": "Strategi mendapatkan, mengonversi, dan mempertahankan customer secara berulang.",
        "formula": "GTM = ICP + channel + pesan utama + funnel + conversion + CAC + retention motion.",
        "example": "Lead dari komunitas UMKM -> webinar -> trial -> onboarding WhatsApp -> paid subscription.",
        "used_in": "Go-To-Market",
    },
    {
        "category": "GTM",
        "term": "Conversion Rate",
        "simple": "Persentase orang/prospek yang berpindah dari satu tahap funnel ke tahap berikutnya.",
        "formula": "Conversion rate = jumlah yang berhasil lanjut / jumlah awal x 100%.",
        "example": "1.000 leads, 120 trial -> conversion lead-to-trial 12%.",
        "used_in": "GTM, Traction",
    },
    {
        "category": "Financials",
        "term": "Operating Cost",
        "simple": "Biaya menjalankan startup, seperti tim, marketing, sales, tools, legal, operasional, dan kantor.",
        "formula": "Operating cost = payroll + marketing + sales + product/engineering + tools + operations + legal/admin.",
        "example": "Payroll Rp 80 juta + marketing Rp 30 juta + tools Rp 10 juta = operating cost Rp 120 juta/bulan.",
        "used_in": "Financials",
    },
    {
        "category": "Financials",
        "term": "Burn Rate / Net Burn",
        "simple": "Jumlah uang yang habis setiap bulan setelah memperhitungkan revenue.",
        "formula": "Net burn = cash keluar bulanan - cash masuk bulanan.",
        "example": "Biaya Rp 120 juta/bulan, revenue Rp 40 juta/bulan -> net burn Rp 80 juta/bulan.",
        "used_in": "Financials, Runway",
    },
    {
        "category": "Financials",
        "term": "EBITDA / Profit",
        "simple": "Indikasi laba operasional sebelum beberapa komponen akuntansi; di aplikasi ini dipakai sebagai estimasi profit/loss sederhana.",
        "formula": "Profit sederhana = revenue - operating cost. EBITDA lebih detail mengecualikan interest, tax, depreciation, amortization.",
        "example": "Revenue Rp 1,2M, cost Rp 2M -> profit/loss Rp -800 juta.",
        "used_in": "Financials",
    },
    {
        "category": "Financials",
        "term": "Profit Margin",
        "simple": "Persentase profit dibanding revenue.",
        "formula": "Profit margin = profit / revenue x 100%.",
        "example": "Profit Rp 3,5M dari revenue Rp 13,5M -> margin 25,9%.",
        "used_in": "Investor Readiness, Financials",
    },
    {
        "category": "Competition",
        "term": "Direct Competitor",
        "simple": "Produk/perusahaan yang menyelesaikan masalah serupa untuk customer yang sama.",
        "formula": "Tidak ada rumus. Bandingkan berdasarkan target customer, use case, pricing, channel, dan workflow.",
        "example": "Aplikasi akuntansi lain yang juga menargetkan UMKM.",
        "used_in": "Competition",
    },
    {
        "category": "Competition",
        "term": "Indirect Competitor",
        "simple": "Alternatif yang menyelesaikan sebagian masalah dengan cara berbeda.",
        "formula": "Tidak ada rumus. Cari alat atau proses yang dipakai customer saat ini walaupun bukan kategori produk yang sama.",
        "example": "Spreadsheet, jasa admin, atau POS tanpa modul analitik.",
        "used_in": "Competition",
    },
    {
        "category": "Competition",
        "term": "Status Quo",
        "simple": "Cara lama yang tetap dipakai customer jika mereka tidak membeli solusi baru.",
        "formula": "Tidak ada rumus. Identifikasi workflow manual, kebiasaan, atau keputusan 'tidak melakukan apa-apa'.",
        "example": "Owner tetap mencatat transaksi di buku tulis karena belum percaya software.",
        "used_in": "Competition",
    },
    {
        "category": "Competition",
        "term": "Moat / Defensibility",
        "simple": "Alasan bisnis sulit ditiru atau dikalahkan ketika sudah tumbuh.",
        "formula": "Bisa berasal dari data, network effect, switching cost, distribution advantage, brand, regulasi, atau operational know-how.",
        "example": "Data transaksi UMKM yang makin akurat membuat rekomendasi produk makin sulit ditiru.",
        "used_in": "Competition, Q&A investor",
    },
]


# ==============================
# Pitch Duration & Business Model System
# ==============================
PITCH_DURATION_OPTIONS = [3, 5, 8, 10, 15, 20]

PITCH_DURATION_PROFILES = {
    3: {"name": "3-minute elevator pitch", "scheme": "7 slide ringkas: konteks tetap lengkap, tetapi beberapa topik digabung agar cocok untuk opening meeting atau demo day singkat.", "slide_goal": "Setiap slide 20-35 detik. Fokus pada satu pesan utama per slide."},
    5: {"name": "5-minute demo day pitch", "scheme": "8 slide kompak: problem, product, market, traction, competition, financial, team, dan ask tetap muncul.", "slide_goal": "Setiap slide 30-45 detik. Gunakan transisi cepat dan hindari membaca semua bullet."},
    8: {"name": "8-minute seed pitch", "scheme": "10-11 slide seed ringkas: cukup detail untuk investor meeting awal tanpa membuat deck terlalu panjang.", "slide_goal": "Setiap slide 35-55 detik. Tekankan bukti dan asumsi utama."},
    10: {"name": "10-minute seed standard", "scheme": "11 slide standar seed: semua konteks investor muncul dengan kompetisi dan milestone dipadatkan.", "slide_goal": "Setiap slide 45-65 detik. Sisakan waktu untuk Q&A."},
    15: {"name": "15-minute detailed seed meeting", "scheme": "Full deck: slide competition dan milestones bisa dipisah/paginated bila data banyak.", "slide_goal": "Setiap slide 50-80 detik. Cocok untuk first investor meeting yang butuh pendalaman."},
    20: {"name": "20-minute deep-dive seed meeting", "scheme": "Full deck dengan ruang narasi lebih longgar. PDF skenario memberi timing lebih detail untuk Q&A dan penjelasan asumsi.", "slide_goal": "Setiap slide 60-100 detik. Gunakan waktu tambahan untuk asumsi, competition, dan milestones."},
}

BUSINESS_MODEL_TEMPLATES = {
    "SaaS / Subscription": {
        "description": "Customer membayar biaya berlangganan bulanan/tahunan untuk memakai software atau platform.",
        "best_for": "Software B2B/B2C, workflow tool, vertical SaaS, productivity, AI tool.",
        "pitch_focus": "Tekankan MRR/ARR, retention, churn, ARPU, CAC payback, dan expansion revenue.",
        "formula": "MRR = pelanggan berbayar x harga bulanan rata-rata; ARR = MRR x 12.",
        "metrics": [("MRR / ARR", "Rp 85 juta MRR"), ("ARPU", "Rp 99.000/bulan"), ("Gross Margin", "78%"), ("CAC Payback", "< 2 bulan")],
        "default_lines": "Subscription SaaS bulanan.\nPaket Pro untuk multi-outlet.\nAdd-on premium untuk laporan, automation, atau analytics.\nExpansion revenue melalui seat, outlet, usage, atau modul tambahan.",
    },
    "Marketplace / Take Rate": {
        "description": "Platform mempertemukan supply dan demand, lalu mengambil komisi dari transaksi.",
        "best_for": "Marketplace jasa, produk, B2B procurement, booking, creator/merchant platform.",
        "pitch_focus": "Tekankan GMV, take rate, liquidity, repeat transaction, supply-demand balance, dan unit economics per transaksi.",
        "formula": "Revenue = GMV x take rate; GMV = jumlah transaksi x nilai transaksi rata-rata.",
        "metrics": [("GMV", "Rp 1,5M/bulan"), ("Take Rate", "5%"), ("Repeat Rate", "42%"), ("Liquidity", "70% order matched")],
        "default_lines": "Marketplace mengambil komisi dari transaksi.\nSupply diakuisisi melalui partner dan komunitas.\nDemand datang dari channel digital dan referral.\nRevenue naik melalui GMV, take rate, dan repeat transaction.",
    },
    "E-commerce / D2C": {
        "description": "Startup menjual produk langsung ke customer dengan margin dari selisih harga jual dan biaya produk/logistik.",
        "best_for": "Brand consumer, D2C, retail online, produk fisik, commerce-enabled brand.",
        "pitch_focus": "Tekankan AOV, gross margin, repeat purchase, CAC, payback, inventory turn, dan kontribusi margin.",
        "formula": "Revenue = jumlah order x AOV; gross margin = (revenue - COGS) / revenue.",
        "metrics": [("AOV", "Rp 180.000"), ("Gross Margin", "45%"), ("Repeat Purchase", "35%"), ("CAC Payback", "< 3 bulan")],
        "default_lines": "Penjualan langsung ke customer melalui website, marketplace, dan channel komunitas.\nMargin berasal dari produk inti dan bundle.\nGrowth ditopang repeat purchase, subscription box, dan channel partnership.",
    },
    "Transaction Fee / Fintech": {
        "description": "Startup mendapat fee dari pembayaran, pembiayaan, transfer, insurance, atau transaksi finansial lain.",
        "best_for": "Fintech, payment, lending enablement, insurtech, embedded finance.",
        "pitch_focus": "Tekankan transaction volume, fee rate, default/risk control bila lending, compliance, CAC, dan contribution margin.",
        "formula": "Revenue = transaction volume x fee rate; untuk lending tambahkan risk cost/default rate.",
        "metrics": [("Transaction Volume", "Rp 3M/bulan"), ("Fee Rate", "1,5%"), ("Default/Risk Cost", "< 2%"), ("Active Merchants", "1.200")],
        "default_lines": "Revenue berasal dari fee transaksi dan revenue share dengan partner finansial.\nProduk tertanam di workflow customer.\nRisk, compliance, dan partner distribution menjadi kunci scale.",
    },
    "Usage-Based / API": {
        "description": "Customer membayar sesuai pemakaian: API call, token, storage, transaksi, seat aktif, atau volume data.",
        "best_for": "Developer tool, AI API, infrastructure, data platform, automation, cloud service.",
        "pitch_focus": "Tekankan usage growth, net revenue retention, gross margin per usage, developer adoption, dan enterprise expansion.",
        "formula": "Revenue = volume usage x harga per unit; NRR penting untuk melihat expansion.",
        "metrics": [("Usage Volume", "2 juta API calls/bulan"), ("Revenue / Unit", "Rp 12/call"), ("Gross Margin", "72%"), ("NRR", "115%")],
        "default_lines": "Customer membayar berdasarkan pemakaian.\nEntry point murah, revenue naik seiring volume.\nExpansion terjadi ketika produk tertanam dalam workflow atau sistem customer.",
    },
    "Freemium": {
        "description": "Sebagian user memakai gratis, lalu sebagian dikonversi ke paket berbayar.",
        "best_for": "Consumer app, productivity, creator tool, education, AI tool, developer product.",
        "pitch_focus": "Tekankan activation, free-to-paid conversion, retention, viral/referral loop, dan cost to serve free users.",
        "formula": "Paid users = free active users x conversion rate; MRR = paid users x ARPU.",
        "metrics": [("Active Free Users", "25.000"), ("Free-to-Paid", "4%"), ("Paid ARPU", "Rp 79.000/bulan"), ("D30 Retention", "38%")],
        "default_lines": "Free plan digunakan untuk akuisisi dan activation.\nRevenue berasal dari upgrade ke paket Pro.\nGrowth ditopang referral, komunitas, dan product-led growth.",
    },
    "Enterprise / Licensing": {
        "description": "Customer perusahaan membayar kontrak tahunan, lisensi, implementation fee, atau seat enterprise.",
        "best_for": "B2B enterprise SaaS, cybersecurity, HR tech, govtech, data/AI enterprise.",
        "pitch_focus": "Tekankan ACV, sales cycle, pipeline, pilot-to-contract conversion, renewal, dan procurement readiness.",
        "formula": "ARR = jumlah kontrak aktif x ACV; pipeline weighted = pipeline value x probability.",
        "metrics": [("ACV", "Rp 250 juta/tahun"), ("Pipeline", "Rp 2,4M"), ("Sales Cycle", "90 hari"), ("Pilot Conversion", "40%")],
        "default_lines": "Revenue berasal dari kontrak tahunan dan implementation fee.\nSales motion berbasis pilot, stakeholder mapping, dan procurement.\nExpansion melalui seat, modul tambahan, dan multi-department rollout.",
    },
    "Service-Enabled Software": {
        "description": "Startup menggabungkan software dengan layanan operasional agar customer cepat mendapat outcome.",
        "best_for": "Managed marketplace, vertical SaaS awal, AI service, operational startup, B2B solution.",
        "pitch_focus": "Tekankan software leverage, margin improvement, repeatability, automation roadmap, dan proses yang bisa distandarisasi.",
        "formula": "Gross margin membaik jika porsi software/automation naik dan biaya layanan per customer turun.",
        "metrics": [("Service Margin", "35%"), ("Automation Rate", "55%"), ("Revenue / Customer", "Rp 2 juta/bulan"), ("Delivery SLA", "< 24 jam")],
        "default_lines": "Layanan membantu customer mendapat outcome cepat.\nSoftware mengurangi biaya delivery dan membuat proses repeatable.\nMargin naik seiring automation dan standardisasi operasi.",
    },
    "Advertising / Media": {
        "description": "Revenue berasal dari iklan, sponsorship, placement, affiliate, atau monetisasi audience.",
        "best_for": "Media, community, content platform, consumer app dengan audience besar.",
        "pitch_focus": "Tekankan audience quality, engagement, CPM/CPC, fill rate, retention, dan diversifikasi revenue.",
        "formula": "Ad revenue = impressions / 1000 x CPM; affiliate revenue = conversion x commission.",
        "metrics": [("MAU", "250.000"), ("Engagement", "12 menit/session"), ("CPM", "Rp 35.000"), ("Fill Rate", "60%")],
        "default_lines": "Revenue berasal dari sponsorship, ads, affiliate, dan brand partnership.\nKekuatan utama adalah audience quality, engagement, dan data segmentasi.\nDiversifikasi revenue dilakukan melalui premium community atau commerce.",
    },
    "Hybrid / Other": {
        "description": "Gabungan beberapa model revenue. Cocok bila startup punya revenue utama dan revenue tambahan.",
        "best_for": "Startup yang masih mencari pricing, multi-product, atau model campuran SaaS + marketplace + services.",
        "pitch_focus": "Tekankan model utama terlebih dahulu. Model tambahan harus memperkuat, bukan membuat cerita membingungkan.",
        "formula": "Total revenue = revenue stream utama + revenue stream tambahan. Pisahkan asumsi per stream.",
        "metrics": [("Primary Revenue", "Subscription"), ("Secondary Revenue", "Revenue share"), ("Gross Margin", "65%"), ("Expansion Driver", "Add-on / usage")],
        "default_lines": "Revenue utama berasal dari model inti.\nRevenue tambahan berasal dari add-on, partnership, atau transaction fee.\nPitch tetap fokus pada satu engine utama agar investor tidak bingung.",
    },
}

STARTUP_TERMS.extend([
    {"category": "Business Model", "term": "SaaS / Subscription", "simple": "Model berlangganan bulanan/tahunan. Investor melihat kualitasnya dari MRR, retention, churn, gross margin, dan CAC payback.", "formula": "MRR = pelanggan berbayar x harga bulanan rata-rata. ARR = MRR x 12.", "example": "1.000 pelanggan x Rp 99.000/bulan = Rp 99 juta MRR.", "used_in": "Business Model, Traction, Financials"},
    {"category": "Business Model", "term": "Marketplace / Take Rate", "simple": "Platform mengambil komisi dari transaksi antara supply dan demand.", "formula": "Revenue = GMV x take rate.", "example": "GMV Rp 1,5M x take rate 5% = Rp 75 juta revenue.", "used_in": "Business Model, Traction"},
    {"category": "Business Model", "term": "Enterprise / Licensing", "simple": "Perusahaan membayar kontrak tahunan, lisensi, atau implementation fee.", "formula": "ARR = jumlah kontrak aktif x ACV.", "example": "10 kontrak x Rp 250 juta/tahun = Rp 2,5M ARR.", "used_in": "Business Model, Financials, GTM"},
    {"category": "Business Model", "term": "Usage-Based Pricing", "simple": "Customer membayar berdasarkan pemakaian seperti API call, token, storage, atau transaksi.", "formula": "Revenue = volume usage x harga per unit.", "example": "2 juta API calls x Rp 12/call = Rp 24 juta revenue.", "used_in": "Business Model, Financials"},
])


def glossary_categories() -> list[str]:
    categories = []
    for item in STARTUP_TERMS:
        category = item.get("category", "Lainnya")
        if category not in categories:
            categories.append(category)
    return categories


def filter_glossary(search: str = "", category: str = "Semua") -> list[dict[str, str]]:
    search_clean = (search or "").strip().lower()
    result = []

    for item in STARTUP_TERMS:
        haystack = " ".join(str(v) for v in item.values()).lower()
        if category != "Semua" and item.get("category") != category:
            continue
        if search_clean and search_clean not in haystack:
            continue
        result.append(item)

    return result


def term_by_name(term_name: str) -> dict[str, str] | None:
    target = term_name.lower().strip()
    for item in STARTUP_TERMS:
        if item["term"].lower().startswith(target) or target in item["term"].lower():
            return item
    return None



# ==============================
# Layout Safety Helpers
# ==============================
def estimate_text_capacity(w: float, h: float, size: float) -> int:
    """Approximate how many characters fit in a textbox before it becomes visually crowded."""
    chars_per_line = max(12, int(w * 9.8 * (12 / max(size, 1))))
    lines_available = max(1, int((h * 72) / (max(size, 1) * 1.25)))
    return max(24, chars_per_line * lines_available)


def truncate_text(text: Any, max_chars: int | None = None) -> str:
    value = str(text or "").strip()
    value = re.sub(r"\s+", " ", value)

    if max_chars is None or len(value) <= max_chars:
        return value

    return value[: max(0, max_chars - 1)].rstrip() + "…"


def adaptive_font_size(text: Any, base_size: int, min_size: int = 8) -> int:
    value = str(text or "")
    length = len(value)

    if length > 220:
        return max(min_size, base_size - 6)
    if length > 150:
        return max(min_size, base_size - 5)
    if length > 95:
        return max(min_size, base_size - 3)
    if length > 55:
        return max(min_size, base_size - 2)

    return base_size


def chunk_items(items: list[Any], chunk_size: int) -> list[list[Any]]:
    return [items[i:i + chunk_size] for i in range(0, len(items), chunk_size)]


def normalize_competitors(data: dict[str, Any]) -> list[dict[str, str]]:
    """Return clean competitor rows from dynamic UI data or legacy fields."""
    rows = data.get("competitors") or []
    cleaned = []

    for row in rows:
        if not isinstance(row, dict):
            continue

        name = str(row.get("name", "")).strip()
        weakness = str(row.get("weakness", "")).strip()
        advantage = str(row.get("advantage", "")).strip()
        category = str(row.get("category", "Alternative")).strip() or "Alternative"

        if name or weakness or advantage:
            cleaned.append(
                {
                    "name": name or "Alternative",
                    "category": category,
                    "weakness": weakness or "Belum diisi",
                    "advantage": advantage or "Belum diisi",
                }
            )

    if cleaned:
        return cleaned

    legacy_rows = [
        {
            "name": data.get("competitor_1", ""),
            "category": "Direct competitor",
            "weakness": data.get("weakness_1", ""),
            "advantage": data.get("advantage_1", ""),
        },
        {
            "name": data.get("competitor_2", ""),
            "category": "Indirect competitor",
            "weakness": data.get("weakness_2", ""),
            "advantage": data.get("advantage_2", ""),
        },
        {
            "name": "Status quo",
            "category": "Status quo",
            "weakness": data.get("status_quo", ""),
            "advantage": data.get("advantage_3", ""),
        },
    ]

    return [row for row in legacy_rows if str(row.get("name", "")).strip()]


def normalize_milestones(data: dict[str, Any]) -> list[dict[str, str]]:
    """Return clean milestone rows from structured UI data or milestone headline."""
    rows = data.get("milestones") or []
    cleaned = []

    for row in rows:
        if not isinstance(row, dict):
            continue

        period = str(row.get("period", "")).strip()
        target = str(row.get("target", "")).strip()
        metric = str(row.get("metric", "")).strip()
        owner = str(row.get("owner", "")).strip()

        if period or target or metric or owner:
            cleaned.append(
                {
                    "period": period or "Next",
                    "target": target or "Target belum diisi",
                    "metric": metric or "Metric belum diisi",
                    "owner": owner or "Team",
                }
            )

    if cleaned:
        return cleaned

    fallback = str(data.get("milestone", "")).strip()

    if fallback:
        return [
            {
                "period": "12-18 bulan",
                "target": fallback,
                "metric": "Funding milestone",
                "owner": "Core team",
            }
        ]

    return []


def milestone_headline(data: dict[str, Any]) -> str:
    milestones = normalize_milestones(data)

    if not milestones:
        return str(data.get("milestone", "Milestone belum diisi"))

    first = milestones[0]
    period = first.get("period", "Next")
    target = first.get("target", "Milestone belum diisi")

    return f"{period}: {target}"



def slugify_key(value: str) -> str:
    value = re.sub(r"[^a-zA-Z0-9]+", "_", str(value or "").lower()).strip("_")
    return value or "default"


def nearest_pitch_duration(value: Any) -> int:
    try:
        minutes = int(value)
    except Exception:
        minutes = 10
    return min(PITCH_DURATION_OPTIONS, key=lambda option: abs(option - minutes))


def pitch_duration_profile(data: dict[str, Any]) -> dict[str, Any]:
    minutes = nearest_pitch_duration(data.get("pitch_duration_minutes", 10))
    profile = dict(PITCH_DURATION_PROFILES.get(minutes, PITCH_DURATION_PROFILES[10]))
    profile["minutes"] = minutes
    if minutes <= 3:
        profile["mode"] = "elevator"
        profile["target_slide_count"] = 7
    elif minutes <= 5:
        profile["mode"] = "demo_day"
        profile["target_slide_count"] = 8
    elif minutes <= 10:
        profile["mode"] = "standard_short"
        profile["target_slide_count"] = 11
    else:
        profile["mode"] = "full"
        profile["target_slide_count"] = "adaptive"
    return profile


def business_model_template(data_or_type: Any) -> dict[str, Any]:
    if isinstance(data_or_type, dict):
        model_type = data_or_type.get("business_model_type", "SaaS / Subscription")
    else:
        model_type = data_or_type
    return BUSINESS_MODEL_TEMPLATES.get(str(model_type), BUSINESS_MODEL_TEMPLATES["SaaS / Subscription"])


def get_business_model_slide_metrics(data: dict[str, Any]) -> list[tuple[str, str]]:
    labels = data.get("model_metric_labels") or []
    values = data.get("model_metric_values") or []
    pairs: list[tuple[str, str]] = []
    for label, value in zip(labels, values):
        label_clean = str(label or "").strip()
        value_clean = str(value or "").strip()
        if label_clean or value_clean:
            pairs.append((label_clean or "Metric", value_clean or "Belum diisi"))
    if pairs:
        return pairs[:4]
    template = business_model_template(data)
    return [(label, value) for label, value in template.get("metrics", [])][:4]


def business_model_key_data(data: dict[str, Any]) -> list[str]:
    template = business_model_template(data)
    model_type = data.get("business_model_type", "SaaS / Subscription")
    result = [f"Model: {model_type}", f"Fokus: {template.get('pitch_focus', '')}"]
    for label, value in get_business_model_slide_metrics(data):
        result.append(f"{label}: {value}")
    result.extend(lines(data.get("business_model", ""), 4))
    return result[:8]


def business_model_formula_text(data: dict[str, Any]) -> str:
    template = business_model_template(data)
    return str(template.get("formula", "Sesuaikan revenue driver dengan model bisnis utama."))


def build_pitch_rhythm(data: dict[str, Any]) -> list[str]:
    profile = pitch_duration_profile(data)
    minutes = profile["minutes"]
    model_type = data.get("business_model_type", "SaaS / Subscription")
    if minutes <= 3:
        return [
            "Opening: 20-30 detik untuk one-liner dan ask.",
            "Problem + Solution + Product: 60-70 detik untuk membuat urgency dan value creation jelas.",
            "Market + Business Model + Traction: 60-70 detik untuk menunjukkan opportunity dan demand.",
            "Competition + Milestones + Ask + Team: 60-70 detik untuk menutup dengan execution plan dan kredibilitas.",
            f"Model bisnis {model_type}: sebutkan 1 revenue driver dan 1 metric paling penting saja.",
        ]
    if minutes <= 5:
        return [
            "Opening: 30 detik untuk company, one-liner, dan ask.",
            "Problem + Product: 90 detik untuk menjelaskan pain, solution, dan demo flow.",
            "Market + Model + Traction: 120 detik untuk membuktikan peluang dan demand awal.",
            "Competition + Financial + Milestones + Team + Ask: 120-150 detik untuk menjawab kesiapan eksekusi.",
            f"Model bisnis {model_type}: tampilkan metric utama dan cara revenue dihitung.",
        ]
    if minutes <= 10:
        return [
            "Opening: 45 detik untuk company, one-liner, round, dan ask.",
            "Problem + Solution + Product: 3 menit untuk membangun urgency dan value creation.",
            "Market + Business Model + Traction + GTM: 4 menit untuk membuktikan opportunity, revenue engine, dan demand.",
            "Competition + Financials + Milestones + Team + Ask: 3 menit untuk menunjukkan positioning, execution plan, dan readiness.",
            f"Model bisnis {model_type}: jelaskan revenue formula dan unit economics paling relevan.",
        ]
    return [
        "Opening: 30-45 detik untuk company, one-liner, dan ask.",
        "Problem + Solution + Product: 3-4 menit untuk membangun urgency dan value creation.",
        "Market + Business Model + Traction + GTM: 4-5 menit untuk membuktikan opportunity dan demand.",
        "Competition + Financials + Milestones: 3-5 menit untuk menjawab positioning dan execution plan.",
        "Team + Fundraising Ask + Closing: 2-3 menit untuk menutup dengan confidence dan next step.",
        f"Model bisnis {model_type}: sesuaikan metrik utama dengan revenue engine, bukan memakai semua metrik startup secara generik.",
    ]


def format_timing(seconds: int) -> str:
    seconds = max(10, int(seconds))
    low = max(10, seconds - 8)
    high = seconds + 8
    if high >= 60:
        return f"{low // 60}:{low % 60:02d}-{high // 60}:{high % 60:02d} menit"
    return f"{low}-{high} detik"


def apply_scenario_timing(items: list[dict[str, Any]], data: dict[str, Any]) -> list[dict[str, Any]]:
    profile = pitch_duration_profile(data)
    total_seconds = profile["minutes"] * 60
    if not items:
        return items
    weights = []
    for item in items:
        title = item.get("title", "").lower()
        weight = 1.0
        if any(keyword in title for keyword in ["problem", "solution", "product", "business", "traction"]):
            weight = 1.25
        if any(keyword in title for keyword in ["financial", "milestone", "ask"]):
            weight = 1.15
        if any(keyword in title for keyword in ["cover", "closing"]):
            weight = 0.65
        weights.append(weight)
    buffer_ratio = 0.18 if profile["minutes"] <= 5 else 0.12
    usable_seconds = int(total_seconds * (1 - buffer_ratio))
    total_weight = sum(weights) or 1
    timed_items = []
    for item, weight in zip(items, weights):
        cloned = dict(item)
        seconds = int(usable_seconds * weight / total_weight)
        cloned["timing"] = format_timing(seconds)
        timed_items.append(cloned)
    return timed_items


# ==============================
# Investor Insight Engine
# ==============================
def generate_investor_insights(data: dict[str, Any]) -> dict[str, Any]:
    rev1 = float(data.get("rev1", 0) or 0)
    rev2 = float(data.get("rev2", 0) or 0)
    rev3 = float(data.get("rev3", 0) or 0)
    cost1 = float(data.get("cost1", 0) or 0)
    cost2 = float(data.get("cost2", 0) or 0)
    cost3 = float(data.get("cost3", 0) or 0)
    profit1 = float(data.get("profit1", 0) or 0)
    profit2 = float(data.get("profit2", 0) or 0)
    profit3 = float(data.get("profit3", 0) or 0)
    ask = float(data.get("ask", 0) or 0)
    runway = int(data.get("runway", 0) or 0)
    currency = data.get("currency", "Rp")
    competitor_count = len(normalize_competitors(data))
    milestone_count = len(normalize_milestones(data))

    growth_y2 = safe_div(rev2 - rev1, rev1) if rev1 else None
    growth_y3 = safe_div(rev3 - rev2, rev2) if rev2 else None
    gross_margin = parse_percent(data.get("gross_margin", ""))
    retention = parse_percent(data.get("retention", ""))
    payback = parse_percent(data.get("payback", ""))

    burn_y1 = max(cost1 - rev1, 0)
    monthly_burn_y1 = burn_y1 / 12 if burn_y1 else 0
    estimated_runway = safe_div(ask, monthly_burn_y1) if monthly_burn_y1 else None
    funding_to_y1_cost = safe_div(ask, cost1) if cost1 else None
    year3_profit_margin = safe_div(profit3, rev3) if rev3 else None
    revenue_multiple_y3 = safe_div(rev3, rev1) if rev1 else None

    completeness_fields = [
        "company", "one_liner", "problem", "problem_evidence", "solution", "value_prop",
        "product_flow", "product_benefit", "market_notes", "business_model", "traction_notes",
        "gtm", "team", "founder_fit", "use_of_funds", "next_round", "milestone",
    ]
    filled = sum(1 for field in completeness_fields if str(data.get(field, "")).strip())
    if competitor_count >= 2:
        filled += 1
    if milestone_count >= 2:
        filled += 1
    completeness = filled / (len(completeness_fields) + 2)

    score = 45
    score += int(completeness * 20)

    if growth_y2 is not None and growth_y2 > 1:
        score += 8
    elif growth_y2 is not None and growth_y2 > 0.3:
        score += 4

    if growth_y3 is not None and growth_y3 > 0.75:
        score += 6
    elif growth_y3 is not None and growth_y3 > 0.25:
        score += 3

    if gross_margin is not None and gross_margin >= 70:
        score += 6
    elif gross_margin is not None and gross_margin >= 50:
        score += 3

    if retention is not None and retention >= 70:
        score += 5
    elif retention is not None and retention >= 50:
        score += 2

    if estimated_runway is not None and estimated_runway >= 15:
        score += 5
    elif estimated_runway is not None and estimated_runway >= 9:
        score += 2

    if year3_profit_margin is not None and year3_profit_margin > 0:
        score += 5

    if competitor_count >= 3:
        score += 3

    if milestone_count >= 3:
        score += 4

    score = min(score, 100)

    strengths = []
    risks = []
    recommendations = []

    if growth_y2 is not None and growth_y2 > 0:
        strengths.append(
            f"Proyeksi revenue naik {pct(growth_y2)} dari Year 1 ke Year 2 dan {pct(growth_y3)} dari Year 2 ke Year 3. Narasikan sebagai momentum pertumbuhan, bukan hanya tabel angka."
        )
    else:
        risks.append(
            "Pertumbuhan revenue belum terlihat kuat dari proyeksi. Investor seed perlu melihat jalur pertumbuhan yang masuk akal dan berulang."
        )

    if gross_margin is not None:
        if gross_margin >= 70:
            strengths.append(
                f"Gross margin {gross_margin:.0f}% memberi sinyal model bisnis berpotensi scalable jika CAC dan retention terjaga."
            )
        else:
            risks.append(
                f"Gross margin {gross_margin:.0f}% perlu dijelaskan: apakah margin akan membaik karena automation, volume, atau perubahan pricing."
            )

    if retention is not None:
        if retention >= 70:
            strengths.append(
                f"Retention {retention:.0f}% dapat dipakai sebagai bukti awal bahwa produk mulai menjadi habit atau workflow penting."
            )
        else:
            risks.append(
                f"Retention {retention:.0f}% perlu diimbangi dengan strategi aktivasi dan customer success agar investor tidak membaca traction sebagai sementara."
            )

    if estimated_runway is not None:
        if estimated_runway >= runway * 1.25:
            recommendations.append(
                f"Runway estimasi dari ask dan burn Year 1 sekitar {estimated_runway:.0f} bulan, lebih tinggi dari input {runway} bulan. Jelaskan buffer atau percepatan hiring agar angka tetap kredibel."
            )
        elif estimated_runway < runway * 0.75:
            risks.append(
                f"Runway input {runway} bulan terlihat agresif dibanding estimasi burn Year 1 sekitar {estimated_runway:.0f} bulan. Perlu validasi ulang use of funds."
            )
        else:
            strengths.append(
                f"Ask pendanaan relatif konsisten dengan burn Year 1 dan runway sekitar {runway} bulan."
            )

    if year3_profit_margin is not None:
        if year3_profit_margin > 0:
            strengths.append(
                f"Profit margin Year 3 sekitar {pct(year3_profit_margin)}. Ini membantu menunjukkan jalur menuju efisiensi, walaupun seed deck tetap harus fokus pada growth."
            )
        else:
            risks.append(
                "Year 3 masih belum menunjukkan profit positif. Jelaskan kenapa loss tersebut disengaja untuk growth dan kapan operating leverage muncul."
            )

    if funding_to_y1_cost is not None:
        recommendations.append(
            f"Funding yang diminta setara {multiple(funding_to_y1_cost)} operating cost Year 1. Hubungkan angka ini dengan hiring, distribusi, dan milestone yang bisa diverifikasi."
        )

    if revenue_multiple_y3 is not None:
        recommendations.append(
            f"Revenue Year 3 adalah {multiple(revenue_multiple_y3)} dibanding Year 1. Tambahkan asumsi utama: jumlah customer, ARPU, churn, dan channel conversion."
        )

    if competitor_count < 3:
        risks.append(
            "Kompetitor/alternatif masih terlalu sedikit. Tambahkan kompetitor langsung, tidak langsung, dan status quo agar positioning terlihat matang."
        )
    else:
        recommendations.append(
            f"Slide kompetisi memuat {competitor_count} alternatif. Gunakan narasi: customer saat ini memakai apa, kelemahannya apa, dan kenapa produk Anda menang."
        )

    if milestone_count < 3:
        risks.append(
            "Milestone masih kurang detail. Seed investor biasanya ingin melihat target 3-4 tahap: product, traction, revenue, dan readiness untuk round berikutnya."
        )
    else:
        recommendations.append(
            f"Gunakan {milestone_count} milestone sebagai jembatan antara ask funding, runway, dan next round story."
        )

    model_type = data.get("business_model_type", "SaaS / Subscription")
    model_template = business_model_template(data)
    recommendations.append(
        f"Model bisnis yang dipilih: {model_type}. Saat pitching, fokuskan slide Business Model pada: {model_template.get('pitch_focus', '')}"
    )

    if len(get_business_model_slide_metrics(data)) < 3:
        risks.append(
            "Metrik khusus model bisnis masih kurang. Tambahkan 3-4 metrik yang sesuai dengan revenue engine, misalnya MRR untuk SaaS, GMV/take rate untuk marketplace, atau ACV/pipeline untuk enterprise."
        )

    if not data.get("problem_evidence", "").strip():
        risks.append(
            "Evidence problem belum kuat. Tambahkan data interview, pilot, waiting list, atau biaya masalah yang dialami customer."
        )

    if not data.get("competition_summary", "").strip():
        risks.append(
            "Slide kompetisi belum punya narrative advantage. Investor biasanya ingin tahu kenapa startup ini bisa menang, bukan hanya siapa pesaingnya."
        )

    if len(str(data.get("one_liner", ""))) > 130:
        recommendations.append(
            "One-liner terlalu panjang. Buat menjadi: untuk siapa, masalah apa, solusi apa, hasil bisnis apa."
        )

    if not strengths:
        strengths.append(
            "Deck sudah memiliki struktur seed standar. Fokus berikutnya adalah memperkuat bukti pasar, angka traction, dan narrative advantage."
        )

    if not risks:
        risks.append(
            "Risiko utama belum terlihat dari input. Tetap siapkan jawaban untuk CAC, churn, defensibility, dan asumsi proyeksi."
        )

    if not recommendations:
        recommendations.append(
            "Tambahkan asumsi di speaker notes: sumber data, metode perhitungan market size, dan alasan channel GTM dipilih."
        )

    headline = (
        f"Deck readiness {score}/100. "
        f"Fokus pitching: buktikan urgency problem, kualitas traction, dan hubungan funding dengan milestone."
    )

    return {
        "score": score,
        "headline": headline,
        "metrics": {
            "growth_y2": pct(growth_y2),
            "growth_y3": pct(growth_y3),
            "burn_y1": short_money(burn_y1, currency),
            "estimated_runway": f"{estimated_runway:.0f} bulan" if estimated_runway else "N/A",
            "year3_profit_margin": pct(year3_profit_margin),
            "funding_to_y1_cost": multiple(funding_to_y1_cost),
        },
        "strengths": strengths[:4],
        "risks": risks[:4],
        "recommendations": recommendations[:4],
    }


# ==============================
# UI Guidance Helpers
# ==============================
def guide(title: str, body: str):
    st.markdown(
        f"""
        <div class="guide-box">
            <strong>{title}</strong><br>
            {body}
        </div>
        """,
        unsafe_allow_html=True,
    )


def render_insight_card(title: str, items: list[str]):
    safe_items = "".join(
        f"<li>{html.escape(str(item))}</li>"
        for item in items
    )

    st.markdown(
        f"""
        <div class="insight-card">
            <h4>{html.escape(title)}</h4>
            <ul>{safe_items}</ul>
        </div>
        """,
        unsafe_allow_html=True,
    )


def show_insights(insights: dict[str, Any]):
    st.subheader("📈 Analisa Investor Readiness")
    st.caption(
        "Analisa ini membaca data yang Anda isi dan memberi gambaran apakah narasi pitching sudah kuat, "
        "bagian mana yang perlu dipertajam, serta insight yang bisa dipakai saat presentasi."
    )

    st.progress(insights["score"] / 100)
    st.markdown(
        f"""
        <div class="readable-panel">
            <p><strong>{html.escape(str(insights['headline']))}</strong></p>
        </div>
        """,
        unsafe_allow_html=True,
    )

    metric_cols = st.columns(6)
    metric_labels = [
        ("Growth Y1→Y2", "growth_y2"),
        ("Growth Y2→Y3", "growth_y3"),
        ("Burn Y1", "burn_y1"),
        ("Runway estimasi", "estimated_runway"),
        ("Margin Y3", "year3_profit_margin"),
        ("Funding/Cost Y1", "funding_to_y1_cost"),
    ]

    for col, (label, key) in zip(metric_cols, metric_labels):
        col.metric(label, insights["metrics"][key])

    c1, c2, c3 = st.columns(3)

    with c1:
        render_insight_card("Strength", insights["strengths"])

    with c2:
        render_insight_card("Risk / Investor Question", insights["risks"])

    with c3:
        render_insight_card("Rekomendasi Pitching", insights["recommendations"])


def render_glossary_section():
    st.subheader("📚 Istilah Startup & Cara Menghitung")
    st.caption(
        "Gunakan bagian ini sebagai mini-kamus saat mengisi deck. Setiap istilah dijelaskan dengan bahasa sederhana, "
        "rumus praktis, contoh angka, dan slide tempat istilah itu biasanya dipakai."
    )

    guide(
        "Cara memakai kamus istilah",
        "Cari istilah yang belum dipahami, baca definisi sederhananya, lalu gunakan rumus dan contoh untuk mengisi field terkait. "
        "Untuk founder baru, prioritaskan memahami TAM/SAM/SOM, MRR, ARPU, CAC, gross margin, retention, runway, burn rate, dan milestone.",
    )

    col1, col2 = st.columns([2, 1])
    with col1:
        search = st.text_input(
            "Cari istilah",
            "",
            placeholder="Contoh: CAC, runway, TAM, retention, MRR",
            help="Ketik istilah, rumus, contoh, atau slide yang ingin dicari.",
        )
    with col2:
        category = st.selectbox(
            "Kategori",
            ["Semua"] + glossary_categories(),
            help="Filter istilah berdasarkan topik agar lebih mudah dipelajari.",
        )

    terms = filter_glossary(search, category)
    st.caption(f"Menampilkan {len(terms)} dari {len(STARTUP_TERMS)} istilah.")

    priority_names = [
        "TAM - Total Addressable Market",
        "SAM - Serviceable Available Market",
        "SOM - Serviceable Obtainable Market",
        "MRR - Monthly Recurring Revenue",
        "ARPU - Average Revenue Per User",
        "CAC - Customer Acquisition Cost",
        "Gross Margin",
        "Retention D7 / D30",
        "Runway",
        "Burn Rate / Net Burn",
        "Milestone",
        "CAC Payback",
    ]

    if not search and category == "Semua":
        st.markdown("### Istilah prioritas untuk pitch deck seed")
        cols = st.columns(3)
        for idx, name in enumerate(priority_names):
            term = next((item for item in STARTUP_TERMS if item["term"] == name), None)
            if not term:
                continue
            with cols[idx % 3]:
                st.markdown(
                    f"""
                    <div class="insight-card">
                        <h4>{html.escape(term['term'])}</h4>
                        <p><strong>Arti:</strong> {html.escape(term['simple'])}</p>
                        <p><strong>Rumus:</strong> {html.escape(term['formula'])}</p>
                    </div>
                    """,
                    unsafe_allow_html=True,
                )

    st.markdown("### Semua istilah")
    if not terms:
        st.warning("Istilah tidak ditemukan. Coba kata kunci lain seperti revenue, market, funding, retention, atau CAC.")
        return

    for item in terms:
        with st.expander(f"{item['term']} — {item['category']}"):
            c1, c2 = st.columns(2)
            with c1:
                st.markdown("**Arti sederhana**")
                st.write(item["simple"])
                st.markdown("**Dipakai di slide**")
                st.write(item["used_in"])
            with c2:
                st.markdown("**Rumus / cara menghitung**")
                st.code(item["formula"], language="text")
                st.markdown("**Contoh**")
                st.info(item["example"])


# ==============================
# Slide Drawing Helpers
# ==============================
def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=18,
    color=None,
    bold=False,
    align=PP_ALIGN.LEFT,
    max_chars=None,
):
    color = color or THEME["ink"]
    size = adaptive_font_size(text, size, 7)

    if max_chars is None:
        max_chars = estimate_text_capacity(w, h, size)

    safe_text = truncate_text(text, max_chars)

    box = slide.shapes.add_textbox(
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )

    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)

    paragraph = tf.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)

    run = paragraph.add_run()
    run.text = safe_text
    run.font.name = FONT_HEAD if bold else FONT_BODY
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color

    return box


def add_seed_header(slide, eyebrow, title, subtitle, data):
    accent = rgb(data["accent_color"])

    add_text(
        slide,
        eyebrow.upper(),
        0.7,
        0.35,
        3.8,
        0.25,
        8,
        accent,
        True,
    )

    add_text(
        slide,
        title,
        0.7,
        0.67,
        10.2,
        0.58,
        28,
        THEME["ink"],
        True,
    )

    add_text(
        slide,
        subtitle,
        0.72,
        1.25,
        11.1,
        0.32,
        10,
        THEME["muted"],
    )

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.7),
        Inches(1.7),
        Inches(11.95),
        Inches(0.01),
    )

    line.fill.solid()
    line.fill.fore_color.rgb = THEME["line"]
    line.line.fill.background()


def add_footer(slide, data, page=None, dark=False):
    color = RGBColor(148, 163, 184) if dark else THEME["muted"]

    add_text(
        slide,
        data.get("company", "Pitch Deck"),
        0.7,
        7.08,
        3.6,
        0.22,
        8,
        color,
    )

    add_text(
        slide,
        DEVELOPER_FOOTER,
        4.25,
        7.08,
        4.8,
        0.22,
        8,
        color,
        align=PP_ALIGN.CENTER,
    )

    if page is not None:
        add_text(
            slide,
            str(page).zfill(2),
            12.0,
            7.08,
            0.65,
            0.22,
            8,
            color,
            align=PP_ALIGN.RIGHT,
        )


def add_card(slide, title, body, x, y, w, h, data):
    accent = rgb(data["accent_color"])

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )

    shape.fill.solid()
    shape.fill.fore_color.rgb = THEME["white"]
    shape.line.color.rgb = THEME["line"]

    add_text(
        slide,
        title.upper(),
        x + 0.22,
        y + 0.18,
        w - 0.44,
        0.24,
        7.5,
        accent,
        True,
        max_chars=44,
    )

    body_size = adaptive_font_size(body, 14, 8)
    add_text(
        slide,
        body,
        x + 0.22,
        y + 0.55,
        w - 0.44,
        h - 0.68,
        body_size,
        THEME["ink"],
        True,
        max_chars=estimate_text_capacity(w - 0.44, h - 0.68, body_size),
    )


def add_big_metric(slide, label, value, x, y, w, data):
    accent = rgb(data["accent_color"])
    value_text = str(value or "")

    value_size = 28
    if len(value_text) > 34:
        value_size = 15
    elif len(value_text) > 22:
        value_size = 19
    elif len(value_text) > 14:
        value_size = 23

    add_text(
        slide,
        value_text,
        x,
        y,
        w,
        0.68,
        value_size,
        accent,
        True,
        max_chars=estimate_text_capacity(w, 0.68, value_size),
    )

    add_text(
        slide,
        label.upper(),
        x,
        y + 0.72,
        w,
        0.22,
        7.5,
        THEME["muted"],
        True,
        max_chars=42,
    )


def add_bullets(slide, items, x, y, w, h, size=18):
    box = slide.shapes.add_textbox(
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )

    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)

    clean_items = items or ["Lengkapi poin utama slide ini."]
    max_items = min(len(clean_items), 5)
    item_size = size

    longest = max((len(str(item)) for item in clean_items[:max_items]), default=0)
    if max_items >= 5 or longest > 150:
        item_size = max(9, size - 5)
    elif longest > 100:
        item_size = max(10, size - 4)
    elif longest > 65:
        item_size = max(11, size - 2)

    per_item_chars = max(38, estimate_text_capacity(w, h / max(max_items, 1), item_size) - 8)

    for i, item in enumerate(clean_items[:max_items]):
        paragraph = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        paragraph.text = truncate_text(item, per_item_chars)
        paragraph.level = 0
        paragraph.space_after = Pt(7 if item_size <= 12 else 10)
        paragraph.font.name = FONT_BODY
        paragraph.font.size = Pt(item_size)
        paragraph.font.color.rgb = THEME["ink"]


def add_takeaway(slide, text, data):
    accent = rgb(data["accent_color"])

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.7),
        Inches(6.18),
        Inches(11.95),
        Inches(0.58),
    )

    shape.fill.solid()
    shape.fill.fore_color.rgb = THEME["blue_soft"]
    shape.line.color.rgb = THEME["blue_line"]

    add_text(
        slide,
        f"Investor takeaway: {truncate_text(text, 155)}",
        0.95,
        6.32,
        11.3,
        0.27,
        9.2,
        accent,
        True,
        max_chars=175,
    )


def add_notes(slide, body):
    try:
        slide.notes_slide.notes_text_frame.text = body
    except Exception:
        pass


def add_table(slide, headers, rows, x, y, w, h, data, column_widths=None):
    accent = rgb(data["accent_color"])
    row_count = len(rows) + 1
    col_count = len(headers)

    table_shape = slide.shapes.add_table(
        row_count,
        col_count,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )

    tbl = table_shape.table

    if column_widths and len(column_widths) == col_count:
        total = sum(column_widths)
        for idx, width_ratio in enumerate(column_widths):
            tbl.columns[idx].width = Inches(w * width_ratio / total)

    header_h = 0.42
    body_h = max(0.42, (h - header_h) / max(len(rows), 1))
    tbl.rows[0].height = Inches(header_h)

    for row_index in range(1, row_count):
        tbl.rows[row_index].height = Inches(body_h)

    font_size = 10
    if len(rows) >= 5:
        font_size = 8
    elif any(len(str(cell)) > 80 for row in rows for cell in row):
        font_size = 8.5
    elif len(rows) >= 4:
        font_size = 9

    for col, value in enumerate(headers):
        cell = tbl.cell(0, col)
        cell.text = truncate_text(value, 38)
        cell.fill.solid()
        cell.fill.fore_color.rgb = accent
        cell.margin_left = Inches(0.06)
        cell.margin_right = Inches(0.06)
        cell.margin_top = Inches(0.03)
        cell.margin_bottom = Inches(0.03)

        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.name = FONT_BODY
            paragraph.font.size = Pt(8.5)
            paragraph.font.bold = True
            paragraph.font.color.rgb = THEME["white"]

    for row_i, row in enumerate(rows, 1):
        for col_i, value in enumerate(row):
            cell = tbl.cell(row_i, col_i)
            max_chars = 95
            if col_i == 0:
                max_chars = 36
            elif col_i == 1 and col_count >= 4:
                max_chars = 46
            cell.text = truncate_text(value, max_chars)
            cell.fill.solid()
            cell.fill.fore_color.rgb = THEME["white"]
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = FONT_BODY
                paragraph.font.size = Pt(font_size)
                paragraph.font.color.rgb = THEME["ink"]
                paragraph.space_after = Pt(0)

    return table_shape


def add_competitor_card(slide, row, x, y, w, h, data):
    accent = rgb(data["accent_color"])

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = THEME["white"]
    shape.line.color.rgb = THEME["line"]

    add_text(slide, row.get("category", "Alternative").upper(), x + 0.22, y + 0.18, w - 0.44, 0.22, 7, accent, True)
    add_text(slide, row.get("name", "Alternative"), x + 0.22, y + 0.43, w - 0.44, 0.32, 14, THEME["ink"], True)
    add_text(slide, "Weakness", x + 0.22, y + 0.86, 1.25, 0.18, 7, THEME["muted"], True)
    add_text(slide, row.get("weakness", "Belum diisi"), x + 0.22, y + 1.08, w / 2 - 0.34, 0.43, 8.5, THEME["ink"])
    add_text(slide, "Our advantage", x + w / 2 + 0.03, y + 0.86, 1.75, 0.18, 7, THEME["muted"], True)
    add_text(slide, row.get("advantage", "Belum diisi"), x + w / 2 + 0.03, y + 1.08, w / 2 - 0.30, 0.43, 8.5, THEME["ink"])


def add_competition_slides(prs, data, page):
    competitors = normalize_competitors(data)

    if not competitors:
        competitors = [
            {
                "name": "Status quo",
                "category": "Status quo",
                "weakness": "Customer tetap memakai cara lama.",
                "advantage": "Produk memberi workflow baru yang lebih cepat dan terukur.",
            }
        ]

    chunk_size = 5
    chunks = chunk_items(competitors, chunk_size)

    for chunk_index, chunk in enumerate(chunks):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide, THEME["bg"])

        suffix = "" if len(chunks) == 1 else f" ({chunk_index + 1}/{len(chunks)})"
        add_seed_header(
            slide,
            "Competition",
            f"We win through focus, speed, and distribution{suffix}",
            "Comparison matrix dibuat adaptif agar tetap rapi meski kompetitor lebih dari satu.",
            data,
        )

        rows = [
            [
                row.get("category", "Alternative"),
                row.get("name", "Alternative"),
                row.get("weakness", "Belum diisi"),
                row.get("advantage", "Belum diisi"),
            ]
            for row in chunk
        ]

        table_height = 3.45 if len(chunk) <= 4 else 3.85
        add_table(
            slide,
            ["Category", "Alternative", "Customer pain / weakness", "Why we win"],
            rows,
            0.78,
            2.02,
            11.78,
            table_height,
            data,
            column_widths=[1.4, 1.8, 3.35, 3.35],
        )

        add_takeaway(
            slide,
            data.get("competition_summary", "Jelaskan kenapa startup ini menang dibanding alternatif yang sudah dipakai customer."),
            data,
        )
        add_footer(slide, data, page)
        add_notes(slide, "Gunakan slide ini untuk menunjukkan positioning. Bila teks panjang, ringkas menjadi customer pain dan why we win yang paling kuat.")
        page += 1

    return page


def add_milestone_card(slide, row, x, y, w, h, data):
    accent = rgb(data["accent_color"])

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = THEME["white"]
    shape.line.color.rgb = THEME["line"]

    # Wide timeline layout; keeps long milestone data readable.
    period_w = 1.75
    metric_x = x + 7.25
    target_w = 4.95
    metric_w = w - 7.55

    add_text(
        slide,
        row.get("period", "Next"),
        x + 0.24,
        y + 0.22,
        period_w,
        0.34,
        12,
        accent,
        True,
        max_chars=28,
    )

    add_text(
        slide,
        row.get("target", "Target belum diisi"),
        x + 2.05,
        y + 0.22,
        target_w,
        0.72,
        13,
        THEME["ink"],
        True,
        max_chars=105,
    )

    add_text(
        slide,
        "Success metric",
        metric_x,
        y + 0.22,
        1.65,
        0.18,
        7,
        THEME["muted"],
        True,
        max_chars=28,
    )

    add_text(
        slide,
        row.get("metric", "Metric belum diisi"),
        metric_x,
        y + 0.47,
        metric_w,
        0.50,
        10,
        THEME["ink"],
        max_chars=88,
    )

    add_text(
        slide,
        f"Owner: {row.get('owner', 'Team')}",
        x + 0.24,
        y + h - 0.32,
        w - 0.48,
        0.18,
        7,
        THEME["muted"],
        max_chars=110,
    )


def add_milestone_slides(prs, data, page):
    milestones = normalize_milestones(data)

    if not milestones:
        return page

    chunk_size = 3
    chunks = chunk_items(milestones, chunk_size)

    for chunk_index, chunk in enumerate(chunks):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide, THEME["bg"])

        suffix = "" if len(chunks) == 1 else f" ({chunk_index + 1}/{len(chunks)})"
        add_seed_header(
            slide,
            "Milestones",
            f"Funding converts into measurable execution milestones{suffix}",
            "Timeline dibuat adaptif agar periode, target, metric, dan owner tetap terbaca.",
            data,
        )

        y_positions = [2.00, 3.35, 4.70]

        for idx, row in enumerate(chunk):
            add_milestone_card(slide, row, 0.85, y_positions[idx], 11.45, 1.14, data)

        add_takeaway(
            slide,
            "Milestone yang baik harus measurable: revenue, users, retention, partnership, atau product release.",
            data,
        )
        add_footer(slide, data, page)
        add_notes(slide, "Gunakan milestone untuk menjelaskan kenapa jumlah funding, runway, dan next round logic saling konsisten.")
        page += 1

    return page


def seed_content_slide(
    prs,
    data,
    page,
    eyebrow,
    title,
    subtitle,
    bullets,
    takeaway,
    metric_items=None,
    side_title=None,
    side_body=None,
    speaker_note="",
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["bg"])

    add_seed_header(
        slide,
        eyebrow,
        title,
        subtitle,
        data,
    )

    if metric_items:
        count = len(metric_items)
        width = 11.5 / max(count, 1)

        for idx, item in enumerate(metric_items):
            label, value = item
            add_big_metric(
                slide,
                label,
                value,
                0.85 + idx * width,
                2.15,
                width - 0.2,
                data,
            )

        add_bullets(
            slide,
            bullets,
            0.85,
            3.55,
            10.8,
            2.2,
            18,
        )

    elif side_title:
        add_bullets(
            slide,
            bullets,
            0.85,
            2.05,
            6.3,
            3.7,
            20,
        )

        add_card(
            slide,
            side_title,
            side_body,
            7.45,
            2.10,
            4.75,
            3.05,
            data,
        )

    else:
        add_bullets(
            slide,
            bullets,
            0.85,
            2.05,
            10.8,
            3.8,
            21,
        )

    add_takeaway(
        slide,
        takeaway,
        data,
    )

    add_footer(
        slide,
        data,
        page,
    )

    add_notes(
        slide,
        speaker_note,
    )

    return slide


def add_insight_slide(prs, data, page):
    insights = generate_investor_insights(data)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["bg"])

    add_seed_header(
        slide,
        "Investor Readiness",
        "What the data says about this pitch",
        "Analisa otomatis untuk membantu founder membaca kekuatan, risiko, dan pesan utama sebelum pitching.",
        data,
    )

    add_big_metric(
        slide,
        "Deck readiness score",
        f"{insights['score']}/100",
        0.85,
        2.05,
        3.2,
        data,
    )

    add_card(
        slide,
        "Growth signal",
        f"Y1→Y2: {insights['metrics']['growth_y2']}\nY2→Y3: {insights['metrics']['growth_y3']}",
        4.2,
        2.0,
        2.55,
        1.45,
        data,
    )

    add_card(
        slide,
        "Runway signal",
        f"Burn Y1: {insights['metrics']['burn_y1']}\nRunway est.: {insights['metrics']['estimated_runway']}",
        7.0,
        2.0,
        2.55,
        1.45,
        data,
    )

    add_card(
        slide,
        "Efficiency signal",
        f"Margin Y3: {insights['metrics']['year3_profit_margin']}\nFunding/Cost Y1: {insights['metrics']['funding_to_y1_cost']}",
        9.8,
        2.0,
        2.55,
        1.45,
        data,
    )

    add_text(
        slide,
        "Key strengths",
        0.85,
        3.8,
        3.3,
        0.32,
        13,
        THEME["green"],
        True,
    )
    add_bullets(slide, insights["strengths"][:3], 0.85, 4.18, 3.55, 1.75, 10)

    add_text(
        slide,
        "Investor questions",
        4.75,
        3.8,
        3.3,
        0.32,
        13,
        THEME["amber"],
        True,
    )
    add_bullets(slide, insights["risks"][:3], 4.75, 4.18, 3.55, 1.75, 10)

    add_text(
        slide,
        "Pitching recommendations",
        8.6,
        3.8,
        3.3,
        0.32,
        13,
        rgb(data["accent_color"]),
        True,
    )
    add_bullets(slide, insights["recommendations"][:3], 8.6, 4.18, 3.55, 1.75, 10)

    add_footer(slide, data, page)
    add_notes(
        slide,
        "Gunakan slide ini sebagai internal readiness view. Bila deck dibagikan ke investor, slide ini bisa dipertahankan atau dihapus sesuai kebutuhan.",
    )



# ==============================
# PDF Pitch Scenario Guide
# ==============================
def pdf_escape(value: Any) -> str:
    """Escape text for ReportLab Paragraph while preserving simple line breaks."""
    safe = html.escape(str(value or "").strip())
    safe = re.sub(r"\s*→\s*", " -> ", safe)
    safe = safe.replace("&quot;", "'")
    return safe.replace("\n", "<br/>")


def pdf_bullets(items: list[Any], limit: int = 5) -> str:
    cleaned = []

    for item in items[:limit]:
        value = truncate_text(str(item or "").strip(), 260)
        if value:
            cleaned.append(f"- {pdf_escape(value)}")

    if not cleaned:
        cleaned = ["- Lengkapi bagian ini sebelum latihan pitching."]

    return "<br/>".join(cleaned)


def scenario_slide_items_full(data: dict[str, Any]) -> list[dict[str, Any]]:
    """Create a scenario guide that mirrors the generated PPTX order."""
    competitors = normalize_competitors(data)
    milestones = normalize_milestones(data)
    insights = generate_investor_insights(data)
    competitor_chunks = chunk_items(competitors, 5) or [[]]
    milestone_chunks = chunk_items(milestones, 3)

    items: list[dict[str, Any]] = []

    def add(
        title: str,
        purpose: str,
        key_data: list[str],
        talk_track: list[str],
        transition: str,
        questions: list[str] | None = None,
        timing: str = "45-60 detik",
    ):
        items.append(
            {
                "title": title,
                "purpose": purpose,
                "key_data": key_data,
                "talk_track": talk_track,
                "transition": transition,
                "questions": questions or [],
                "timing": timing,
            }
        )

    add(
        "Cover",
        "Membuka konteks: siapa startupnya, apa yang dilakukan, tahap pendanaan, dan siapa presenter.",
        [
            f"Company: {data.get('company', '')}",
            f"One-liner: {data.get('one_liner', '')}",
            f"Round: {data.get('round', '')}",
            f"Ask: {money(data.get('ask', 0), data.get('currency', 'Rp'))}",
        ],
        [
            "Buka dengan satu kalimat tajam: customer, pain, solusi, dan hasil bisnis.",
            "Jangan mulai dari sejarah panjang perusahaan. Mulai dari kenapa masalah ini penting sekarang.",
        ],
        "Masuk ke problem agar investor memahami urgency sebelum melihat produk.",
        ["Apakah one-liner langsung menjawab customer dan masalah utama?"],
        "30-45 detik",
    )

    add(
        "Problem",
        "Membuktikan bahwa masalah customer nyata, sering terjadi, mahal, dan cukup besar untuk didanai.",
        lines(data.get("problem", ""), 5) + [f"Evidence: {data.get('problem_evidence', '')}"],
        [
            "Ceritakan kondisi customer sebelum memakai produk.",
            "Tekankan frekuensi masalah, biaya masalah, dan konsekuensi jika masalah tidak diselesaikan.",
            "Gunakan evidence singkat: interview, pilot, waiting list, data internal, atau observasi market.",
        ],
        "Setelah pain terasa jelas, lanjutkan ke solusi yang langsung menghilangkan pain tersebut.",
        [
            "Seberapa sering masalah ini terjadi?",
            "Berapa biaya ekonomi atau operasional dari masalah ini?",
            "Kenapa solusi lama belum cukup?",
        ],
    )

    add(
        "Solution",
        "Menunjukkan cara baru yang lebih sederhana, lebih cepat, atau lebih efektif untuk menyelesaikan problem.",
        lines(data.get("solution", ""), 5) + [f"Value proposition: {data.get('value_prop', '')}"],
        [
            "Jelaskan perubahan before-after: dari workflow lama ke workflow baru.",
            "Fokus pada outcome, bukan daftar fitur teknis.",
            "Pastikan value proposition terdengar spesifik dan terukur.",
        ],
        "Lanjutkan ke Product untuk menunjukkan bagaimana solusi tersebut benar-benar bekerja.",
        [
            "Apa yang membuat solusi ini 10x lebih baik?",
            "Bagian mana dari solusi yang paling sulit ditiru?",
        ],
    )

    add(
        "Product",
        "Memperlihatkan mekanisme produk sehingga investor paham flow dan value creation dalam waktu singkat.",
        [
            f"Product flow: {data.get('product_flow', '')}",
            f"Benefit: {data.get('product_benefit', '')}",
        ] + lines(data.get("features", ""), 5),
        [
            "Jelaskan flow: input -> proses -> output -> dampak untuk customer.",
            "Jika ada mockup, arahkan investor ke bagian yang menunjukkan value paling kuat.",
            "Hindari menjelaskan semua tombol. Pilih satu use case utama.",
        ],
        "Setelah produk dipahami, tunjukkan bahwa market-nya cukup besar dan entry wedge-nya jelas.",
        [
            "Apa core workflow yang membuat user kembali memakai produk?",
            "Apakah produk sudah dipakai customer nyata atau masih prototype?",
        ],
        "60-90 detik",
    )

    add(
        "Market",
        "Menunjukkan peluang pasar besar, tetapi tetap dengan fokus segmen awal yang realistis dimenangkan.",
        [
            f"TAM: {data.get('tam', '')}",
            f"SAM: {data.get('sam', '')}",
            f"SOM: {data.get('som', '')}",
        ] + lines(data.get("market_notes", ""), 5),
        [
            "Jelaskan definisi TAM, SAM, dan SOM secara sederhana.",
            "Tekankan wedge awal: customer mana yang paling sakit dan paling cepat diakuisisi.",
            "Sebutkan potensi ekspansi setelah segmen awal terbukti.",
        ],
        "Setelah pasar terlihat menarik, jelaskan bagaimana startup menghasilkan revenue dari pasar tersebut.",
        [
            "Sumber perhitungan market size dari mana?",
            "Kenapa segmen awal ini paling tepat?",
        ],
    )

    add(
        "Business Model",
        "Menjelaskan cara startup menghasilkan uang dan apakah unit economics punya potensi scalable sesuai model bisnis yang dipilih.",
        business_model_key_data(data) + [f"Formula: {business_model_formula_text(data)}"],
        [
            "Jelaskan siapa yang membayar, berapa, kapan, dan kenapa mereka bersedia membayar.",
            "Gunakan metrik yang sesuai dengan model bisnis, bukan metrik generik.",
            "Sebutkan potensi expansion revenue jika ada.",
        ],
        "Setelah model revenue jelas, tunjukkan bukti bahwa market mulai merespons.",
        [
            "Apakah pricing sudah diuji?",
            "Apa revenue driver utama model bisnis ini?",
            "Bagaimana margin membaik saat scale?",
        ],
    )

    add(
        "Traction",
        "Membuktikan demand awal: pengguna, revenue, growth, retention, pilot, LOI, pipeline, atau partnership.",
        [
            f"Users/customers: {data.get('users', '')}",
            f"Revenue/GMV: {data.get('revenue', '')}",
            f"Growth: {data.get('growth', '')}",
            f"Retention: {data.get('retention', '')}",
        ] + lines(data.get("traction_notes", ""), 5),
        [
            "Mulai dari metrik paling kuat: revenue, retention, active usage, atau paid pilot.",
            "Bedakan traction nyata dari vanity metrics.",
            "Jika growth masih awal, jelaskan leading indicators yang kredibel.",
        ],
        "Setelah demand terbukti, jelaskan bagaimana customer acquisition akan dibuat berulang.",
        [
            "Apa metrik traction paling investor-grade?",
            "Retention dihitung dari definisi apa?",
            "Apakah growth berasal dari channel yang repeatable?",
        ],
    )

    add(
        "Go-To-Market",
        "Menjelaskan mesin akuisisi customer: ICP, channel utama, conversion motion, dan alasan channel bisa scale.",
        [
            f"ICP: {data.get('icp', '')}",
            f"Primary channel: {data.get('channel', '')}",
        ] + lines(data.get("gtm", ""), 5),
        [
            "Jelaskan siapa customer pertama yang paling tepat diburu.",
            "Tunjukkan channel yang sudah diuji atau paling masuk akal berdasarkan perilaku customer.",
            "Hubungkan GTM dengan CAC dan payback.",
        ],
        "Setelah channel jelas, bandingkan dengan alternatif yang sudah dipakai customer hari ini.",
        [
            "Channel mana yang sudah terbukti?",
            "Apa bottleneck GTM terbesar?",
            "Bagaimana sales cycle dan conversion rate?",
        ],
    )

    for index, chunk in enumerate(competitor_chunks, 1):
        key_data = []
        for row in chunk:
            key_data.append(
                f"{row.get('name', 'Alternative')} ({row.get('category', 'Alternative')}): weakness = {row.get('weakness', '')}; why we win = {row.get('advantage', '')}"
            )

        title = "Competition" if len(competitor_chunks) == 1 else f"Competition {index}/{len(competitor_chunks)}"
        add(
            title,
            "Menunjukkan bahwa founder memahami medan kompetisi dan punya alasan kuat untuk menang.",
            key_data + [f"Narrative advantage: {data.get('competition_summary', '')}"],
            [
                "Jangan mengatakan tidak ada kompetitor. Status quo juga kompetitor.",
                "Jelaskan customer sekarang memakai apa, kenapa itu kurang, dan kenapa produk Anda lebih tepat.",
                "Gunakan why we win yang spesifik: distribusi, workflow, data, speed, cost, atau domain expertise.",
            ],
            "Setelah positioning jelas, masuk ke financials untuk menunjukkan rencana scale dan penggunaan dana.",
            [
                "Apa moat atau unfair advantage?",
                "Bagaimana defensibility jika kompetitor besar masuk?",
                "Apakah positioning Anda cukup tajam untuk segmen awal?",
            ],
        )

    add(
        "Financials",
        "Menghubungkan proyeksi revenue, cost, profit, runway, dan milestone agar ask funding terlihat logis.",
        [
            f"Revenue Y1/Y2/Y3: {money(data.get('rev1', 0), data.get('currency', 'Rp'))} / {money(data.get('rev2', 0), data.get('currency', 'Rp'))} / {money(data.get('rev3', 0), data.get('currency', 'Rp'))}",
            f"Cost Y1/Y2/Y3: {money(data.get('cost1', 0), data.get('currency', 'Rp'))} / {money(data.get('cost2', 0), data.get('currency', 'Rp'))} / {money(data.get('cost3', 0), data.get('currency', 'Rp'))}",
            f"Profit Y1/Y2/Y3: {money(data.get('profit1', 0), data.get('currency', 'Rp'))} / {money(data.get('profit2', 0), data.get('currency', 'Rp'))} / {money(data.get('profit3', 0), data.get('currency', 'Rp'))}",
            f"Runway: {data.get('runway', '')} bulan",
            f"Next milestone: {milestone_headline(data)}",
        ],
        [
            "Jangan hanya membacakan tabel. Jelaskan 2-3 asumsi utama yang menggerakkan revenue dan cost.",
            "Hubungkan ask funding dengan runway dan milestone berikutnya.",
            "Akui bahwa proyeksi seed adalah asumsi, lalu jelaskan asumsi mana yang akan divalidasi.",
        ],
        "Setelah angka utama jelas, uraikan milestone eksekusi yang akan dicapai dengan pendanaan.",
        [
            "Apa asumsi revenue paling sensitif?",
            "Berapa burn bulanan?",
            "Kapan perlu raise round berikutnya?",
        ],
        "60-90 detik",
    )

    for index, chunk in enumerate(milestone_chunks, 1):
        title = "Milestones" if len(milestone_chunks) == 1 else f"Milestones {index}/{len(milestone_chunks)}"
        add(
            title,
            "Membuat pendanaan terasa konkret: dana digunakan untuk mencapai target yang measurable.",
            [
                f"{row.get('period', 'Next')}: {row.get('target', '')} | Metric: {row.get('metric', '')} | Owner: {row.get('owner', '')}"
                for row in chunk
            ],
            [
                "Baca milestone sebagai timeline eksekusi, bukan daftar keinginan.",
                "Setiap milestone harus punya metric keberhasilan dan owner.",
                "Tekankan milestone yang membuka peluang next round atau inflection point berikutnya.",
            ],
            "Setelah milestone, gunakan investor readiness sebagai cek internal sebelum menutup dengan team dan ask.",
            [
                "Milestone mana yang paling berisiko?",
                "Apa leading indicator sebelum milestone tercapai?",
            ],
        )

    if data.get("include_insight_slide", True):
        add(
            "Investor Readiness",
            "Slide internal/opsional untuk memahami kekuatan, risiko, dan rekomendasi dari data deck.",
            [
                insights.get("headline", ""),
                f"Growth Y1-Y2: {insights['metrics']['growth_y2']}",
                f"Estimated runway: {insights['metrics']['estimated_runway']}",
                f"Margin Y3: {insights['metrics']['year3_profit_margin']}",
            ],
            [
                "Gunakan bagian ini untuk latihan internal sebelum meeting investor.",
                "Ambil 1-2 risiko terbesar dan siapkan jawaban singkat berbasis data.",
                "Slide ini boleh dihapus dari deck investor jika ingin deck lebih ringkas.",
            ],
            "Setelah insight internal, masuk ke Team untuk menunjukkan kemampuan eksekusi.",
            insights.get("risks", []),
        )

    add(
        "Team",
        "Meyakinkan investor bahwa tim punya founder-market fit dan kemampuan eksekusi untuk menang.",
        lines(data.get("team", ""), 5) + [f"Founder-market fit: {data.get('founder_fit', '')}"],
        [
            "Hubungkan pengalaman tim dengan problem, produk, dan market.",
            "Tunjukkan unfair advantage: domain expertise, akses customer, technical edge, atau distribusi.",
            "Jika ada gap tim, jelaskan hiring plan dari funding.",
        ],
        "Setelah tim dipercaya, tutup dengan ask yang jelas dan next milestone.",
        [
            "Kenapa tim ini yang paling tepat?",
            "Role kunci apa yang masih perlu direkrut?",
        ],
    )

    add(
        "Fundraising Ask",
        "Menutup dengan jumlah dana, penggunaan dana, runway, milestone, dan logic menuju round berikutnya.",
        [
            f"Ask: {money(data.get('ask', 0), data.get('currency', 'Rp'))} {data.get('round', '')}",
            f"Runway: {data.get('runway', '')} bulan",
            f"Next milestone: {milestone_headline(data)}",
        ] + lines(data.get("use_of_funds", ""), 5) + [f"Next round logic: {data.get('next_round', '')}"],
        [
            "Nyatakan ask dengan percaya diri dan spesifik.",
            "Jelaskan penggunaan dana dalam 3-4 kategori utama.",
            "Hubungkan funding ke milestone yang membuat startup layak raise round berikutnya.",
        ],
        "Akhiri dengan closing line, ajakan follow-up, demo, data room, atau meeting berikutnya.",
        [
            "Apakah ask cukup untuk 12-18 bulan?",
            "Milestone apa yang membuat round berikutnya lebih kuat?",
            "Apa penggunaan dana terbesar dan kenapa?",
        ],
        "60 detik",
    )

    add(
        "Closing",
        "Mengakhiri pitch dengan visi besar dan ajakan tindak lanjut yang jelas.",
        [f"Closing line: {data.get('closing', '')}", f"Contact: {data.get('contact', '')}"],
        [
            "Rangkum visi besar dalam satu kalimat.",
            "Sampaikan ajakan jelas: demo, follow-up call, data room, atau term discussion.",
            "Berhenti dengan percaya diri dan buka sesi tanya jawab.",
        ],
        "Q&A: siapkan jawaban untuk market size, CAC, retention, competition, runway, dan assumptions.",
        [
            "Apa satu pesan yang harus investor ingat setelah meeting?",
            "Apa next step yang diminta dari investor?",
        ],
        "30 detik",
    )

    return items



def scenario_slide_items_compact(data: dict[str, Any]) -> list[dict[str, Any]]:
    """Scenario guide for duration-adjusted compact deck versions."""
    profile = pitch_duration_profile(data)
    insights = generate_investor_insights(data)
    competitors = normalize_competitors(data)[:3]
    milestones = normalize_milestones(data)[:3]

    def comp_lines() -> list[str]:
        if not competitors:
            return ["Status quo: jelaskan cara lama yang masih dipakai customer dan kenapa produk Anda menang."]
        return [f"{row.get('name', 'Alternative')} - {row.get('weakness', '')}; why we win: {row.get('advantage', '')}" for row in competitors]

    def milestone_lines() -> list[str]:
        if not milestones:
            return [f"Next milestone: {milestone_headline(data)}"]
        return [f"{row.get('period', 'Next')}: {row.get('target', '')} | metric: {row.get('metric', '')}" for row in milestones]

    items: list[dict[str, Any]] = []

    def add(title, purpose, key_data, talk_track, transition, questions=None):
        items.append({"title": title, "purpose": purpose, "key_data": key_data, "talk_track": talk_track, "transition": transition, "questions": questions or [], "timing": "otomatis"})

    add(
        "Cover",
        "Membuka pitch dengan company, one-liner, round, ask, durasi, dan model bisnis.",
        [f"Company: {data.get('company', '')}", f"One-liner: {data.get('one_liner', '')}", f"Ask: {money(data.get('ask', 0), data.get('currency', 'Rp'))} {data.get('round', '')}", f"Pitch duration: {profile['minutes']} menit", f"Business model: {data.get('business_model_type', 'SaaS / Subscription')}",],
        ["Buka dengan kalimat yang langsung menjawab customer, problem, solusi, dan outcome.", "Jangan membaca cover terlalu lama; gunakan sebagai anchor narasi."],
        "Masuk ke problem agar investor memahami urgency sebelum melihat produk.",
        ["Apakah one-liner bisa diucapkan dalam 15 detik?"],
    )

    if profile["minutes"] <= 3:
        add("Problem + Solution + Product", "Menggabungkan pain, solusi, dan product flow agar konteks tetap lengkap dalam durasi sangat singkat.", lines(data.get("problem", ""), 3) + lines(data.get("solution", ""), 3) + [f"Product flow: {data.get('product_flow', '')}"], ["Ceritakan pain paling besar, bukan semua masalah.", "Langsung tunjukkan bagaimana solusi mengubah workflow customer.", "Gunakan product flow input -> proses -> output -> dampak."], "Setelah value creation jelas, tunjukkan market dan cara bisnis menghasilkan revenue.", ["Apa problem paling mahal?", "Apa workflow utama produk?"])
        add("Market + Business Model", "Menjelaskan peluang pasar dan revenue engine tanpa kehilangan konteks model bisnis.", [f"TAM/SAM/SOM: {data.get('tam', '')} / {data.get('sam', '')} / {data.get('som', '')}"] + business_model_key_data(data), ["Sebutkan segmen awal yang paling realistis dimenangkan.", "Jelaskan formula revenue sesuai model bisnis.", "Pilih satu metrik bisnis yang paling kuat untuk disebutkan."], "Setelah market dan model jelas, buktikan demand dan channel akuisisi.", ["Kenapa segmen awal ini?", "Bagaimana revenue dihitung?"])
        add("Traction + Go-To-Market", "Membuktikan demand awal dan cara memperoleh customer secara repeatable.", [f"Users: {data.get('users', '')}", f"Revenue/GMV: {data.get('revenue', '')}", f"Growth: {data.get('growth', '')}", f"Retention: {data.get('retention', '')}"] + lines(data.get("gtm", ""), 3), ["Buka dengan traction paling kuat.", "Hubungkan traction dengan channel utama dan ICP."], "Setelah demand terlihat, jawab positioning dan rencana eksekusi.", ["Traction mana yang paid/active?", "Channel mana yang sudah terbukti?"])
        add("Competition + Milestones", "Menunjukkan positioning dan target eksekusi utama yang akan dicapai dengan funding.", comp_lines() + milestone_lines(), ["Akui alternatif yang dipakai customer hari ini.", "Tunjukkan why we win secara spesifik.", "Sebutkan milestone paling penting untuk 12-18 bulan."], "Tutup dengan team dan ask agar investor tahu siapa yang mengeksekusi dan dana dipakai untuk apa.", ["Apa moat?", "Milestone mana yang paling menentukan next round?"])
        add("Team + Fundraising Ask", "Menggabungkan kredibilitas tim, penggunaan dana, runway, dan next milestone.", lines(data.get("team", ""), 3) + [f"Ask: {money(data.get('ask', 0), data.get('currency', 'Rp'))}", f"Runway: {data.get('runway', '')} bulan", f"Use of funds: {truncate_text(data.get('use_of_funds', ''), 160)}"], ["Hubungkan pengalaman tim dengan market.", "Nyatakan ask dan use of funds secara jelas."], "Akhiri dengan closing line dan next step.", ["Kenapa tim ini tepat?", "Apakah ask cukup untuk runway?"])
    elif profile["minutes"] <= 5:
        add("Problem + Solution", "Membangun urgency dan langsung menunjukkan perubahan before-after.", lines(data.get("problem", ""), 4) + lines(data.get("solution", ""), 4), ["Gunakan problem paling tajam.", "Jelaskan solution sebagai outcome, bukan fitur."], "Lanjutkan ke product agar investor melihat mekanismenya.", ["Apa bukti problem nyata?"])
        add("Product", "Menjelaskan flow produk dan benefit utama.", [f"Product flow: {data.get('product_flow', '')}", f"Benefit: {data.get('product_benefit', '')}"] + lines(data.get("features", ""), 3), ["Jelaskan satu use case utama.", "Tampilkan dampak, bukan semua fitur."], "Setelah product jelas, masuk ke market dan revenue engine.", ["Apa workflow yang membuat user kembali?"])
        add("Market + Business Model", "Menggabungkan market size dan cara startup menghasilkan uang.", [f"TAM/SAM/SOM: {data.get('tam', '')} / {data.get('sam', '')} / {data.get('som', '')}"] + business_model_key_data(data), ["Jelaskan wedge market.", "Gunakan formula revenue sesuai model bisnis."], "Setelah opportunity jelas, tampilkan demand awal.", ["Bagaimana TAM/SAM/SOM dihitung?", "Apa metric utama model bisnis ini?"])
        add("Traction + GTM", "Membuktikan demand dan menunjukkan channel akuisisi yang bisa diulang.", [f"Users: {data.get('users', '')}", f"Revenue/GMV: {data.get('revenue', '')}", f"Growth: {data.get('growth', '')}", f"Retention: {data.get('retention', '')}"] + lines(data.get("gtm", ""), 3), ["Mulai dari metric terkuat.", "Hubungkan GTM dengan ICP dan CAC."], "Setelah demand dan GTM, jawab kompetisi dan rencana milestone.", ["Apa traction paling investor-grade?"])
        add("Competition + Milestones", "Menyatukan positioning dan execution plan agar investor melihat alasan menang dan target berikutnya.", comp_lines() + milestone_lines(), ["Sebutkan alternatif utama.", "Tekankan why we win dan milestone measurable."], "Lanjutkan ke financial dan ask agar funding terasa logis.", ["Apa defensibility?", "Milestone mana yang unlock next round?"])
        add("Financials + Ask", "Menghubungkan proyeksi, runway, use of funds, dan ask.", [f"Revenue Y1/Y2/Y3: {money(data.get('rev1', 0), data.get('currency', 'Rp'))} / {money(data.get('rev2', 0), data.get('currency', 'Rp'))} / {money(data.get('rev3', 0), data.get('currency', 'Rp'))}", f"Runway: {data.get('runway', '')} bulan", f"Ask: {money(data.get('ask', 0), data.get('currency', 'Rp'))}"] + lines(data.get("use_of_funds", ""), 3), ["Jelaskan asumsi utama, bukan membaca semua angka.", "Hubungkan ask ke milestone."], "Tutup dengan team dan closing.", ["Berapa burn bulanan?", "Apa asumsi revenue paling sensitif?"])
        add("Team + Closing", "Menutup dengan founder-market fit dan ajakan tindak lanjut.", lines(data.get("team", ""), 4) + [f"Closing: {data.get('closing', '')}"], ["Tunjukkan kenapa tim tepat.", "Minta next step yang jelas."], "Q&A.", ["Apa gap tim terbesar?"])
    else:
        for title, purpose, key_data, talk_track, transition, questions in [
            ("Problem", "Membuktikan pain yang nyata dan urgent.", lines(data.get("problem", ""), 5) + [f"Evidence: {data.get('problem_evidence', '')}"], ["Tunjukkan frekuensi dan biaya masalah.", "Gunakan bukti problem."], "Lanjut ke solution.", ["Apa bukti problem nyata?"]),
            ("Solution", "Menjelaskan cara baru yang mengubah kondisi customer.", lines(data.get("solution", ""), 5) + [f"Value prop: {data.get('value_prop', '')}"], ["Fokus pada outcome before-after.", "Jangan terlalu teknis."], "Lanjut ke product.", ["Apa yang 10x lebih baik?"]),
            ("Product", "Menunjukkan workflow produk dan benefit utama.", [f"Flow: {data.get('product_flow', '')}", f"Benefit: {data.get('product_benefit', '')}"] + lines(data.get("features", ""), 4), ["Jelaskan input -> proses -> output -> dampak.", "Pilih use case terpenting."], "Lanjut ke market.", ["Apakah sudah dipakai customer?"]),
            ("Market", "Menunjukkan peluang besar dengan wedge awal yang fokus.", [f"TAM: {data.get('tam', '')}", f"SAM: {data.get('sam', '')}", f"SOM: {data.get('som', '')}"] + lines(data.get("market_notes", ""), 4), ["Jelaskan cara hitung dan segmen awal.", "Tekankan wedge market."], "Lanjut ke business model.", ["Bagaimana market size dihitung?"]),
            ("Business Model", "Menjelaskan revenue engine sesuai model bisnis.", business_model_key_data(data) + [f"Formula: {business_model_formula_text(data)}"], ["Jelaskan siapa membayar, berapa, kapan.", "Gunakan metrik sesuai model."], "Lanjut ke traction.", ["Apa revenue driver utama?"]),
            ("Traction", "Membuktikan demand awal.", [f"Users: {data.get('users', '')}", f"Revenue/GMV: {data.get('revenue', '')}", f"Growth: {data.get('growth', '')}", f"Retention: {data.get('retention', '')}"] + lines(data.get("traction_notes", ""), 4), ["Mulai dari traction terkuat.", "Jelaskan definisi metric."], "Lanjut ke GTM dan competition.", ["Apakah traction paid atau vanity?"]),
            ("GTM + Competition", "Menunjukkan cara akuisisi customer dan alasan menang dibanding alternatif.", [f"ICP: {data.get('icp', '')}", f"Channel: {data.get('channel', '')}"] + lines(data.get("gtm", ""), 3) + comp_lines(), ["Hubungkan ICP, channel, dan CAC.", "Sebutkan why we win."], "Lanjut ke financials dan milestones.", ["Apa channel yang repeatable?", "Apa moat?"]),
            ("Financials + Milestones", "Menghubungkan proyeksi, runway, dan target eksekusi.", [f"Revenue Y1/Y2/Y3: {money(data.get('rev1', 0), data.get('currency', 'Rp'))} / {money(data.get('rev2', 0), data.get('currency', 'Rp'))} / {money(data.get('rev3', 0), data.get('currency', 'Rp'))}", f"Runway: {data.get('runway', '')} bulan"] + milestone_lines(), ["Jelaskan 2-3 asumsi utama.", "Tunjukkan milestone measurable."], "Lanjut ke team.", ["Berapa burn?", "Milestone apa yang unlock next round?"]),
            ("Team", "Membuktikan founder-market fit.", lines(data.get("team", ""), 5) + [f"Founder-market fit: {data.get('founder_fit', '')}"], ["Hubungkan pengalaman tim dengan market.", "Sebutkan gap hiring jika ada."], "Tutup dengan ask.", ["Kenapa tim ini tepat?"]),
            ("Fundraising Ask + Closing", "Menutup dengan ask, use of funds, next round logic, dan next step.", [f"Ask: {money(data.get('ask', 0), data.get('currency', 'Rp'))}", f"Runway: {data.get('runway', '')} bulan"] + lines(data.get("use_of_funds", ""), 4) + [f"Closing: {data.get('closing', '')}"], ["Nyatakan ask secara spesifik.", "Hubungkan dana dengan milestone."], "Q&A.", ["Apa penggunaan dana terbesar?", "Kapan next round?"]),
        ]:
            add(title, purpose, key_data, talk_track, transition, questions)
    return apply_scenario_timing(items, data)


def scenario_slide_items(data: dict[str, Any]) -> list[dict[str, Any]]:
    profile = pitch_duration_profile(data)
    if profile["mode"] == "full":
        return apply_scenario_timing(scenario_slide_items_full(data), data)
    return scenario_slide_items_compact(data)


def build_scenario_pdf(data: dict[str, Any]) -> BytesIO:
    """Build a PDF speaker guide that teaches the pitch scenario order."""
    output = BytesIO()
    doc = SimpleDocTemplate(
        output,
        pagesize=A4,
        rightMargin=1.45 * cm,
        leftMargin=1.45 * cm,
        topMargin=1.3 * cm,
        bottomMargin=1.55 * cm,
        title=f"{data.get('company', 'Startup')} - Pitch Scenario Guide",
        author=DEVELOPER_FOOTER,
    )

    base_styles = getSampleStyleSheet()
    styles = {
        "cover_title": ParagraphStyle(
            "CoverTitle",
            parent=base_styles["Title"],
            fontName="Helvetica-Bold",
            fontSize=24,
            leading=30,
            textColor=colors.HexColor("#0f172a"),
            spaceAfter=14,
        ),
        "cover_subtitle": ParagraphStyle(
            "CoverSubtitle",
            parent=base_styles["BodyText"],
            fontName="Helvetica",
            fontSize=11,
            leading=16,
            textColor=colors.HexColor("#475569"),
            spaceAfter=10,
        ),
        "section_title": ParagraphStyle(
            "SectionTitle",
            parent=base_styles["Heading2"],
            fontName="Helvetica-Bold",
            fontSize=15,
            leading=19,
            textColor=colors.HexColor("#0f172a"),
            spaceBefore=4,
            spaceAfter=8,
        ),
        "label": ParagraphStyle(
            "Label",
            parent=base_styles["BodyText"],
            fontName="Helvetica-Bold",
            fontSize=8,
            leading=10,
            textColor=colors.HexColor("#2563eb"),
            spaceAfter=2,
        ),
        "body": ParagraphStyle(
            "Body",
            parent=base_styles["BodyText"],
            fontName="Helvetica",
            fontSize=9.2,
            leading=13.2,
            textColor=colors.HexColor("#1e293b"),
            spaceAfter=6,
        ),
        "small": ParagraphStyle(
            "Small",
            parent=base_styles["BodyText"],
            fontName="Helvetica",
            fontSize=8,
            leading=11,
            textColor=colors.HexColor("#64748b"),
        ),
        "center": ParagraphStyle(
            "Center",
            parent=base_styles["BodyText"],
            fontName="Helvetica",
            fontSize=9,
            leading=12,
            textColor=colors.HexColor("#64748b"),
            alignment=TA_CENTER,
        ),
    }

    story = []
    company = data.get("company", "Startup")
    one_liner = data.get("one_liner", "")
    insights = generate_investor_insights(data)
    scenario_items = scenario_slide_items(data)

    def add_footer_canvas(canvas, doc_obj):
        canvas.saveState()
        canvas.setFont("Helvetica", 7.5)
        canvas.setFillColor(colors.HexColor("#64748b"))
        canvas.drawString(1.45 * cm, 0.85 * cm, DEVELOPER_FOOTER)
        canvas.drawRightString(A4[0] - 1.45 * cm, 0.85 * cm, f"Page {doc_obj.page}")
        canvas.restoreState()

    story.append(Paragraph(pdf_escape(f"{company} - Pitch Scenario Guide"), styles["cover_title"]))
    story.append(Paragraph(pdf_escape(one_liner), styles["cover_subtitle"]))
    story.append(Spacer(1, 0.25 * cm))

    profile = pitch_duration_profile(data)
    summary_rows = [
        ["Purpose", "PDF latihan pitching berdasarkan urutan slide PPTX yang di-generate."],
        ["Durasi pitching", f"{profile['minutes']} menit - {profile['name']}"],
        ["Skema deck", profile["scheme"]],
        ["Business model", f"{data.get('business_model_type', 'SaaS / Subscription')} - {business_model_template(data).get('description', '')}"],
        ["Round", str(data.get("round", ""))],
        ["Ask", money(data.get("ask", 0), data.get("currency", "Rp"))],
        ["Readiness", f"{insights['score']}/100"],
        ["How to use", "Latih setiap slide sesuai timing, pahami transisi, lalu siapkan jawaban untuk pertanyaan investor."],
    ]
    summary_table = RLTable(
        [[Paragraph(pdf_escape(a), styles["label"]), Paragraph(pdf_escape(b), styles["body"])] for a, b in summary_rows],
        colWidths=[3.2 * cm, 12.7 * cm],
        hAlign="LEFT",
    )
    summary_table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, -1), colors.HexColor("#f8fafc")),
                ("BOX", (0, 0), (-1, -1), 0.6, colors.HexColor("#cbd5e1")),
                ("INNERGRID", (0, 0), (-1, -1), 0.35, colors.HexColor("#e2e8f0")),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("LEFTPADDING", (0, 0), (-1, -1), 8),
                ("RIGHTPADDING", (0, 0), (-1, -1), 8),
                ("TOPPADDING", (0, 0), (-1, -1), 6),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
            ]
        )
    )
    story.append(summary_table)
    story.append(Spacer(1, 0.45 * cm))

    story.append(Paragraph("Recommended pitch rhythm", styles["section_title"]))
    rhythm = build_pitch_rhythm(data)
    story.append(Paragraph(pdf_bullets(rhythm, 8), styles["body"]))
    story.append(PageBreak())

    story.append(Paragraph("Slide-by-slide scenario", styles["cover_title"]))
    story.append(Paragraph("Ikuti urutan ini saat latihan. Setiap bagian berisi tujuan slide, data yang harus disebutkan, narasi, transisi, dan pertanyaan investor yang perlu disiapkan.", styles["cover_subtitle"]))
    story.append(Spacer(1, 0.15 * cm))

    for idx, item in enumerate(scenario_items, 1):
        story.append(Paragraph(pdf_escape(f"{idx:02d}. {item['title']}"), styles["section_title"]))

        meta = RLTable(
            [
                [Paragraph("Timing", styles["label"]), Paragraph(pdf_escape(item.get("timing", "45-60 detik")), styles["body"])],
                [Paragraph("Purpose", styles["label"]), Paragraph(pdf_escape(item.get("purpose", "")), styles["body"])],
                [Paragraph("Key data to mention", styles["label"]), Paragraph(pdf_bullets(item.get("key_data", []), 7), styles["body"])],
                [Paragraph("Talk track", styles["label"]), Paragraph(pdf_bullets(item.get("talk_track", []), 5), styles["body"])],
                [Paragraph("Transition", styles["label"]), Paragraph(pdf_escape(item.get("transition", "")), styles["body"])],
                [Paragraph("Investor questions", styles["label"]), Paragraph(pdf_bullets(item.get("questions", []), 5), styles["body"])],
            ],
            colWidths=[3.4 * cm, 12.5 * cm],
            hAlign="LEFT",
        )
        meta.setStyle(
            TableStyle(
                [
                    ("BOX", (0, 0), (-1, -1), 0.5, colors.HexColor("#cbd5e1")),
                    ("INNERGRID", (0, 0), (-1, -1), 0.3, colors.HexColor("#e2e8f0")),
                    ("BACKGROUND", (0, 0), (0, -1), colors.HexColor("#eff6ff")),
                    ("BACKGROUND", (1, 0), (1, -1), colors.white),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ("LEFTPADDING", (0, 0), (-1, -1), 7),
                    ("RIGHTPADDING", (0, 0), (-1, -1), 7),
                    ("TOPPADDING", (0, 0), (-1, -1), 5),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
                ]
            )
        )
        story.append(meta)
        story.append(Spacer(1, 0.28 * cm))

        if idx in {5, 10, 15} and idx != len(scenario_items):
            story.append(PageBreak())

    story.append(PageBreak())
    story.append(Paragraph("Model bisnis startup dan metrik utama", styles["cover_title"]))
    story.append(Paragraph("Bagian ini membantu founder memilih model bisnis yang paling cocok. Saat pitching, jangan menampilkan semua metrik; pilih metrik yang sesuai dengan revenue engine utama.", styles["cover_subtitle"]))
    story.append(Spacer(1, 0.15 * cm))

    selected_model = data.get("business_model_type", "SaaS / Subscription")
    for model_name, model in BUSINESS_MODEL_TEMPLATES.items():
        title = model_name if model_name != selected_model else f"{model_name} - model yang dipilih"
        metric_text = "; ".join([f"{label}: {value}" for label, value in model.get("metrics", [])])
        model_table = RLTable(
            [
                [Paragraph("Model", styles["label"]), Paragraph(pdf_escape(title), styles["body"])],
                [Paragraph("Cocok untuk", styles["label"]), Paragraph(pdf_escape(model.get("best_for", "")), styles["body"])],
                [Paragraph("Cara pitch", styles["label"]), Paragraph(pdf_escape(model.get("pitch_focus", "")), styles["body"])],
                [Paragraph("Rumus", styles["label"]), Paragraph(pdf_escape(model.get("formula", "")), styles["body"])],
                [Paragraph("Metrik", styles["label"]), Paragraph(pdf_escape(metric_text), styles["body"])],
            ],
            colWidths=[3.4 * cm, 12.5 * cm],
            hAlign="LEFT",
        )
        model_table.setStyle(TableStyle([
            ("BOX", (0, 0), (-1, -1), 0.45, colors.HexColor("#cbd5e1")),
            ("INNERGRID", (0, 0), (-1, -1), 0.25, colors.HexColor("#e2e8f0")),
            ("BACKGROUND", (0, 0), (0, -1), colors.HexColor("#eff6ff") if model_name == selected_model else colors.HexColor("#f8fafc")),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("LEFTPADDING", (0, 0), (-1, -1), 6),
            ("RIGHTPADDING", (0, 0), (-1, -1), 6),
            ("TOPPADDING", (0, 0), (-1, -1), 4),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
        ]))
        story.append(model_table)
        story.append(Spacer(1, 0.16 * cm))

    story.append(PageBreak())
    story.append(Paragraph("Istilah investor dan cara menghitungnya", styles["cover_title"]))
    story.append(
        Paragraph(
            "Appendix ini membantu founder baru memahami istilah yang muncul di form, slide, dan Q&A investor. "
            "Gunakan rumus sebagai pendekatan praktis; sesuaikan dengan jenis bisnis dan kualitas data yang tersedia.",
            styles["cover_subtitle"],
        )
    )
    story.append(Spacer(1, 0.15 * cm))

    active_category = None
    for item in STARTUP_TERMS:
        if item.get("category") != active_category:
            active_category = item.get("category")
            story.append(Spacer(1, 0.18 * cm))
            story.append(Paragraph(pdf_escape(active_category), styles["section_title"]))

        term_table = RLTable(
            [
                [Paragraph("Istilah", styles["label"]), Paragraph(pdf_escape(item.get("term", "")), styles["body"])],
                [Paragraph("Arti sederhana", styles["label"]), Paragraph(pdf_escape(item.get("simple", "")), styles["body"])],
                [Paragraph("Rumus / cara hitung", styles["label"]), Paragraph(pdf_escape(item.get("formula", "")), styles["body"])],
                [Paragraph("Contoh", styles["label"]), Paragraph(pdf_escape(item.get("example", "")), styles["body"])],
                [Paragraph("Dipakai di slide", styles["label"]), Paragraph(pdf_escape(item.get("used_in", "")), styles["body"])],
            ],
            colWidths=[3.4 * cm, 12.5 * cm],
            hAlign="LEFT",
        )
        term_table.setStyle(
            TableStyle(
                [
                    ("BOX", (0, 0), (-1, -1), 0.45, colors.HexColor("#cbd5e1")),
                    ("INNERGRID", (0, 0), (-1, -1), 0.25, colors.HexColor("#e2e8f0")),
                    ("BACKGROUND", (0, 0), (0, -1), colors.HexColor("#f8fafc")),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ("LEFTPADDING", (0, 0), (-1, -1), 6),
                    ("RIGHTPADDING", (0, 0), (-1, -1), 6),
                    ("TOPPADDING", (0, 0), (-1, -1), 4),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
                ]
            )
        )
        story.append(term_table)
        story.append(Spacer(1, 0.16 * cm))

    story.append(PageBreak())
    story.append(Paragraph("Final rehearsal checklist", styles["cover_title"]))
    checklist = [
        "One-liner dapat diucapkan dalam kurang dari 15 detik.",
        "Problem memiliki bukti, bukan hanya opini.",
        "Product demo hanya menampilkan workflow yang paling penting.",
        "Traction memakai metrik demand, bukan vanity metrics.",
        "Competition mencakup kompetitor langsung, tidak langsung, dan status quo.",
        "Financials memiliki asumsi yang bisa dijelaskan.",
        "Milestone terhubung langsung dengan jumlah funding dan runway.",
        "Ask dinyatakan jelas: jumlah, penggunaan dana, runway, dan next milestone.",
        "Q&A sudah disiapkan untuk CAC, retention, market size, defensibility, dan burn rate.",
    ]
    story.append(Paragraph(pdf_bullets(checklist, 12), styles["body"]))
    story.append(Spacer(1, 0.35 * cm))
    story.append(Paragraph(pdf_escape(DEVELOPER_FOOTER), styles["center"]))

    doc.build(story, onFirstPage=add_footer_canvas, onLaterPages=add_footer_canvas)
    output.seek(0)
    return output


# ==============================
# Deck Generator
# ==============================
def build_full_deck(data, image_buffer=None):
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    accent = rgb(data["accent_color"])
    page = 1

    # 1. Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["dark"])

    add_text(slide, data["company"], 0.8, 0.85, 8.5, 0.65, 36, THEME["white"], True)
    add_text(slide, data["one_liner"], 0.82, 1.75, 8.9, 0.85, 24, RGBColor(226, 232, 240), True)
    add_text(slide, f'{data["round"]} • {money(data["ask"], data["currency"])}', 0.85, 4.85, 7.5, 0.35, 15, accent, True)
    add_text(slide, f'{data["presenter"]} | {data["contact"]}', 0.85, 5.35, 7.5, 0.28, 11, RGBColor(203, 213, 225))
    add_text(slide, "Seed Investor Pitch", 9.6, 6.72, 2.8, 0.3, 10, RGBColor(148, 163, 184), True, PP_ALIGN.RIGHT)

    add_footer(slide, data, page, dark=True)
    add_notes(slide, "Open with a concise narrative: customer, pain, solution, and why this can become venture-scale.")
    page += 1

    seed_content_slide(
        prs,
        data,
        page,
        "Problem",
        "The current workflow is broken and expensive",
        "Tunjukkan pain point yang sering terjadi, mahal, dan cukup besar untuk menjadi venture-scale opportunity.",
        lines(data["problem"], 5),
        "Problem harus terasa urgent, bukan sekadar nice-to-have.",
        side_title="Proof of pain",
        side_body=data["problem_evidence"],
        speaker_note="Mulai dari pain point, frekuensi masalah, dan biaya masalah. Jangan mulai dari fitur.",
    )
    page += 1

    seed_content_slide(
        prs,
        data,
        page,
        "Solution",
        "A simpler way to solve the problem",
        "Jelaskan perubahan sebelum dan sesudah produk dipakai.",
        lines(data["solution"], 5),
        "Solusi harus langsung menjawab problem utama.",
        side_title="Core value proposition",
        side_body=data["value_prop"],
        speaker_note="Jelaskan perubahan kondisi pelanggan sebelum dan sesudah memakai produk.",
    )
    page += 1

    # Product
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["bg"])
    add_seed_header(slide, "Product", "The product turns user activity into business outcomes", "Gunakan slide ini untuk demo flow utama, bukan daftar fitur panjang.", data)

    if image_buffer:
        image_buffer.seek(0)
        slide.shapes.add_picture(image_buffer, Inches(0.85), Inches(2.05), width=Inches(6.2), height=Inches(3.45))
        add_card(slide, "Product flow", data["product_flow"], 7.45, 2.05, 4.75, 1.75, data)
        add_card(slide, "Product benefit", data["product_benefit"], 7.45, 4.10, 4.75, 1.35, data)
    else:
        add_bullets(slide, lines(data["features"], 5), 0.85, 2.05, 6.4, 3.7, 20)
        add_card(slide, "Product flow", data["product_flow"], 7.45, 2.15, 4.75, 1.85, data)
        add_card(slide, "Product benefit", data["product_benefit"], 7.45, 4.25, 4.75, 1.25, data)

    add_takeaway(slide, "Investor harus paham cara produk menciptakan nilai dalam 30 detik.", data)
    add_footer(slide, data, page)
    add_notes(slide, "Demo singkat: input, proses, output, dan dampak untuk pelanggan.")
    page += 1

    seed_content_slide(
        prs,
        data,
        page,
        "Market",
        "A large market with a focused entry wedge",
        "TAM menunjukkan potensi; SOM menunjukkan fokus eksekusi awal.",
        lines(data["market_notes"], 5),
        "Pasar awal harus spesifik dan bisa dimenangkan.",
        metric_items=[("TAM", data["tam"]), ("SAM", data["sam"]), ("SOM", data["som"])],
        speaker_note="Tekankan wedge pasar awal, bukan hanya ukuran pasar yang besar.",
    )
    page += 1

    seed_content_slide(
        prs,
        data,
        page,
        "Business Model",
        "Revenue can scale with attractive unit economics",
        "Tampilkan cara menghasilkan uang dan indikasi kualitas bisnis.",
        lines(data["business_model"], 5),
        "Seed investor mencari sinyal bahwa bisnis bisa menjadi scalable.",
        metric_items=get_business_model_slide_metrics(data),
        speaker_note="Jelaskan siapa yang membayar, kapan membayar, dan kenapa margin bisa naik saat skala naik.",
    )
    page += 1

    seed_content_slide(
        prs,
        data,
        page,
        "Traction",
        "Early demand is already visible",
        "Untuk seed, traction bisa berupa revenue, active usage, retention, pilot, LOI, atau pipeline.",
        lines(data["traction_notes"], 5),
        "Traction harus membuktikan demand, bukan vanity metrics.",
        metric_items=[("Users", data["users"]), ("Revenue", data["revenue"]), ("Growth", data["growth"]), ("Retention", data["retention"])],
        speaker_note="Gunakan metrik demand: revenue, active usage, retention, pilot, LOI, atau pipeline.",
    )
    page += 1

    seed_content_slide(
        prs,
        data,
        page,
        "Go-To-Market",
        "A repeatable acquisition motion is emerging",
        "Jelaskan ICP, channel utama, dan cara distribusi bisa berulang.",
        lines(data["gtm"], 5),
        "GTM harus menunjukkan mesin akuisisi, bukan hanya rencana marketing.",
        side_title="ICP & channel",
        side_body=f'{data["icp"]}\n\nPrimary channel: {data["channel"]}',
        speaker_note="Tunjukkan lead source, conversion, dan biaya akuisisi pelanggan.",
    )
    page += 1

    page = add_competition_slides(prs, data, page)

    # Financials
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["bg"])
    add_seed_header(slide, "Financials", "Funding converts into measurable milestones", "Seed financials harus sederhana, berbasis asumsi, dan terhubung ke runway.", data)
    add_table(
        slide,
        ["Metric", "Year 1", "Year 2", "Year 3"],
        [
            ["Revenue", money(data["rev1"], data["currency"]), money(data["rev2"], data["currency"]), money(data["rev3"], data["currency"])],
            ["Operating Cost", money(data["cost1"], data["currency"]), money(data["cost2"], data["currency"]), money(data["cost3"], data["currency"])],
            ["EBITDA / Profit", money(data["profit1"], data["currency"]), money(data["profit2"], data["currency"]), money(data["profit3"], data["currency"])],
        ],
        0.85,
        2.05,
        11.65,
        2.7,
        data,
    )
    add_big_metric(slide, "Runway", f'{data["runway"]} months', 0.95, 5.3, 3.0, data)
    add_big_metric(slide, "Next milestone", milestone_headline(data), 4.25, 5.3, 7.4, data)
    add_footer(slide, data, page)
    add_notes(slide, "Jelaskan asumsi utama. Hubungkan funding ke milestone berikutnya.")
    page += 1

    page = add_milestone_slides(prs, data, page)

    if data.get("include_insight_slide", True):
        add_insight_slide(prs, data, page)
        page += 1

    seed_content_slide(
        prs,
        data,
        page,
        "Team",
        "The right team for this market",
        "Di seed stage, investor banyak menilai founder-market fit.",
        lines(data["team"], 5),
        "Tim harus terlihat punya unfair advantage untuk mengeksekusi bisnis ini.",
        side_title="Founder-market fit",
        side_body=data["founder_fit"],
        speaker_note="Hubungkan pengalaman tim dengan problem dan akses ke market.",
    )
    page += 1

    seed_content_slide(
        prs,
        data,
        page,
        "Fundraising Ask",
        "We are raising to reach the next inflection point",
        "Hubungkan jumlah dana dengan runway, penggunaan dana, dan milestone berikutnya.",
        lines(data["use_of_funds"], 5),
        "Ask harus jelas: jumlah, penggunaan, runway, dan target pembuktian.",
        side_title=data["round"],
        side_body=f'{money(data["ask"], data["currency"])}\n\nNext milestone: {milestone_headline(data)}\n\n{data["next_round"]}',
        speaker_note="Tutup dengan jumlah dana, runway, use of funds, dan target 12-18 bulan.",
    )
    page += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["dark"])
    add_text(slide, data["closing"], 0.85, 1.25, 9.8, 1.1, 34, THEME["white"], True)
    add_text(slide, f'{data["presenter"]}\n{data["contact"]}', 0.9, 5.35, 7.0, 0.6, 14, RGBColor(203, 213, 225))
    add_footer(slide, data, page, dark=True)
    add_notes(slide, "Akhiri dengan ajakan follow-up: demo, data room, atau meeting berikutnya.")

    output = BytesIO()
    prs.save(output)
    output.seek(0)

    return output



def add_compact_split_slide(
    prs,
    data,
    page,
    eyebrow,
    title,
    subtitle,
    left_title,
    left_items,
    right_title,
    right_items,
    takeaway,
    metric_items=None,
    speaker_note="",
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["bg"])
    add_seed_header(slide, eyebrow, title, subtitle, data)
    if metric_items:
        count = len(metric_items)
        width = 11.5 / max(count, 1)
        for idx, (label, value) in enumerate(metric_items):
            add_big_metric(slide, label, value, 0.85 + idx * width, 2.05, width - 0.15, data)
        content_y = 3.35
        content_h = 2.25
    else:
        content_y = 2.15
        content_h = 3.55
    add_text(slide, left_title, 0.85, content_y, 5.4, 0.32, 13, rgb(data["accent_color"]), True, max_chars=60)
    add_bullets(slide, left_items, 0.85, content_y + 0.42, 5.55, content_h - 0.28, 12.5)
    add_text(slide, right_title, 6.85, content_y, 5.4, 0.32, 13, rgb(data["accent_color"]), True, max_chars=60)
    add_bullets(slide, right_items, 6.85, content_y + 0.42, 5.45, content_h - 0.28, 12.5)
    add_takeaway(slide, takeaway, data)
    add_footer(slide, data, page)
    add_notes(slide, speaker_note)
    return slide


def add_compact_cover(prs, data, page):
    profile = pitch_duration_profile(data)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["dark"])
    add_text(slide, data["company"], 0.8, 0.8, 8.6, 0.65, 35, THEME["white"], True)
    add_text(slide, data["one_liner"], 0.82, 1.72, 9.4, 0.85, 23, RGBColor(226, 232, 240), True)
    add_text(slide, f'{data["round"]} • {money(data["ask"], data["currency"])}', 0.85, 4.65, 7.8, 0.35, 15, rgb(data["accent_color"]), True)
    add_text(slide, f'Duration: {profile["minutes"]} min • {data.get("business_model_type", "SaaS / Subscription")}', 0.85, 5.12, 8.8, 0.28, 11, RGBColor(203, 213, 225))
    add_text(slide, f'{data["presenter"]} | {data["contact"]}', 0.85, 5.48, 7.5, 0.28, 11, RGBColor(203, 213, 225))
    add_text(slide, profile["name"], 9.35, 6.72, 3.0, 0.3, 10, RGBColor(148, 163, 184), True, PP_ALIGN.RIGHT)
    add_footer(slide, data, page, dark=True)
    add_notes(slide, "Buka dengan satu kalimat tajam: customer, problem, solusi, outcome, dan ask.")


def add_compact_closing(prs, data, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, THEME["dark"])
    add_text(slide, data.get("closing", "Thank you"), 0.85, 1.15, 9.8, 1.1, 33, THEME["white"], True)
    add_text(slide, f'{data.get("presenter", "")}\n{data.get("contact", "")}', 0.9, 5.35, 7.0, 0.6, 14, RGBColor(203, 213, 225))
    add_text(slide, "Next step: demo, data room, follow-up call, or term discussion", 0.9, 4.70, 8.8, 0.35, 12, RGBColor(203, 213, 225))
    add_footer(slide, data, page, dark=True)
    add_notes(slide, "Akhiri dengan ajakan follow-up yang spesifik.")


def compact_competition_lines(data: dict[str, Any], limit: int = 3) -> list[str]:
    competitors = normalize_competitors(data)[:limit]
    if not competitors:
        return ["Status quo: jelaskan cara lama customer dan kenapa produk Anda menang."]
    return [f"{row.get('name', 'Alternative')}: {row.get('advantage', '')}" for row in competitors]


def compact_milestone_lines(data: dict[str, Any], limit: int = 3) -> list[str]:
    milestones = normalize_milestones(data)[:limit]
    if not milestones:
        return [milestone_headline(data)]
    return [f"{row.get('period', 'Next')}: {row.get('target', '')} ({row.get('metric', '')})" for row in milestones]


def build_compact_deck(data, image_buffer=None):
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    profile = pitch_duration_profile(data)
    page = 1
    add_compact_cover(prs, data, page)
    page += 1
    if profile["minutes"] <= 3:
        add_compact_split_slide(prs, data, page, "Problem / Solution / Product", "Pain, solution, and product flow in one clear story", "Durasi sangat singkat: gabungkan konteks utama tanpa menghilangkan alur investor.", "Problem + Solution", lines(data.get("problem", ""), 3) + lines(data.get("solution", ""), 3), "Product flow", [data.get("product_flow", ""), data.get("product_benefit", "")] + lines(data.get("features", ""), 2), "Investor harus memahami pain dan value creation sebelum melihat angka.")
        page += 1
        add_compact_split_slide(prs, data, page, "Market / Business Model", "A focused market with a clear revenue engine", "Market size dan model bisnis dipadatkan agar investor tetap melihat opportunity dan cara monetisasi.", "Market wedge", [f"TAM: {data.get('tam', '')}", f"SAM: {data.get('sam', '')}", f"SOM: {data.get('som', '')}"] + lines(data.get("market_notes", ""), 2), "Revenue engine", business_model_key_data(data), "Jangan hanya menyebut pasar besar; jelaskan wedge awal dan cara revenue dihitung.", metric_items=get_business_model_slide_metrics(data))
        page += 1
        add_compact_split_slide(prs, data, page, "Traction / GTM", "Early demand and repeatable acquisition motion", "Gabungkan bukti demand dan channel agar investor melihat momentum serta cara scale.", "Traction", [f"Users: {data.get('users', '')}", f"Revenue/GMV: {data.get('revenue', '')}", f"Growth: {data.get('growth', '')}", f"Retention: {data.get('retention', '')}"] + lines(data.get("traction_notes", ""), 2), "Go-to-market", [f"ICP: {data.get('icp', '')}", f"Channel: {data.get('channel', '')}"] + lines(data.get("gtm", ""), 3), "Traction harus membuktikan demand, GTM harus menunjukkan cara demand diulang.")
        page += 1
        add_compact_split_slide(prs, data, page, "Competition / Milestones", "Why we win and what this funding unlocks", "Positioning dan milestone digabung agar konteks kompetisi dan eksekusi tetap ada.", "Competition snapshot", compact_competition_lines(data, 3) + [data.get("competition_summary", "")], "Execution milestones", compact_milestone_lines(data, 3), "Investor perlu melihat alasan menang dan target eksekusi yang measurable.")
        page += 1
        add_compact_split_slide(prs, data, page, "Team / Ask", "The right team raising to reach the next inflection point", "Tutup dengan founder-market fit, ask, runway, use of funds, dan milestone berikutnya.", "Team", lines(data.get("team", ""), 4) + [data.get("founder_fit", "")], "Fundraising ask", [f"Ask: {money(data.get('ask', 0), data.get('currency', 'Rp'))}", f"Runway: {data.get('runway', '')} bulan", f"Next milestone: {milestone_headline(data)}"] + lines(data.get("use_of_funds", ""), 3), "Ask harus terasa sebagai jembatan menuju milestone, bukan sekadar kebutuhan kas.")
        page += 1
        add_compact_closing(prs, data, page)
    elif profile["minutes"] <= 5:
        add_compact_split_slide(prs, data, page, "Problem / Solution", "The current workflow is broken, and the fix is clear", "Gabungkan pain dan solusi untuk menjaga tempo demo day.", "Problem", lines(data.get("problem", ""), 4) + [f"Evidence: {data.get('problem_evidence', '')}"], "Solution", lines(data.get("solution", ""), 4) + [data.get("value_prop", "")], "Problem harus terasa urgent dan solution harus langsung menjawabnya.")
        page += 1
        add_compact_split_slide(prs, data, page, "Product", "The product creates value through a simple workflow", "Pilih satu use case utama agar investor paham product value dalam 45 detik.", "Product flow", [data.get("product_flow", ""), data.get("product_benefit", "")], "Key features", lines(data.get("features", ""), 5), "Demo flow harus menunjukkan input, proses, output, dan dampak.")
        page += 1
        add_compact_split_slide(prs, data, page, "Market / Business Model", "A focused wedge with a model-specific revenue engine", "Model bisnis disesuaikan dengan jenis revenue yang dipilih.", "Market", [f"TAM: {data.get('tam', '')}", f"SAM: {data.get('sam', '')}", f"SOM: {data.get('som', '')}"] + lines(data.get("market_notes", ""), 2), "Business model", business_model_key_data(data), "Pasar awal harus spesifik; revenue engine harus mudah dihitung.", metric_items=get_business_model_slide_metrics(data))
        page += 1
        add_compact_split_slide(prs, data, page, "Traction / GTM", "Early demand is visible and acquisition can repeat", "Tampilkan metric demand dan channel utama.", "Traction", [f"Users: {data.get('users', '')}", f"Revenue/GMV: {data.get('revenue', '')}", f"Growth: {data.get('growth', '')}", f"Retention: {data.get('retention', '')}"], "GTM", [f"ICP: {data.get('icp', '')}", f"Channel: {data.get('channel', '')}"] + lines(data.get("gtm", ""), 3), "Gabungkan demand signal dengan repeatable acquisition motion.")
        page += 1
        add_compact_split_slide(prs, data, page, "Competition / Milestones", "Positioning and execution plan", "Jelaskan why we win dan target berikutnya secara ringkas.", "Competition", compact_competition_lines(data, 3) + [data.get("competition_summary", "")], "Milestones", compact_milestone_lines(data, 3), "Why we win dan milestones harus measurable.")
        page += 1
        add_compact_split_slide(prs, data, page, "Financials / Ask", "Funding converts into measurable milestones", "Hubungkan proyeksi, runway, use of funds, dan ask.", "Financials", [f"Revenue Y1: {money(data.get('rev1', 0), data.get('currency', 'Rp'))}", f"Revenue Y2: {money(data.get('rev2', 0), data.get('currency', 'Rp'))}", f"Revenue Y3: {money(data.get('rev3', 0), data.get('currency', 'Rp'))}", f"Runway: {data.get('runway', '')} bulan"], "Ask", [f"Ask: {money(data.get('ask', 0), data.get('currency', 'Rp'))}", f"Next milestone: {milestone_headline(data)}"] + lines(data.get("use_of_funds", ""), 3), "Angka harus menjelaskan why now, why this amount, dan what it unlocks.")
        page += 1
        add_compact_split_slide(prs, data, page, "Team / Closing", "The right team for this market", "Tutup dengan founder-market fit dan next step.", "Team", lines(data.get("team", ""), 4) + [data.get("founder_fit", "")], "Closing", [data.get("closing", ""), f"Contact: {data.get('contact', '')}"], "Investor harus ingat masalah, alasan menang, dan next step.")
    else:
        seed_content_slide(prs, data, page, "Problem", "The current workflow is broken and expensive", "Tunjukkan pain point yang sering terjadi, mahal, dan cukup besar untuk menjadi venture-scale opportunity.", lines(data.get("problem", ""), 5), "Problem harus terasa urgent, bukan sekadar nice-to-have.", side_title="Proof of pain", side_body=data.get("problem_evidence", ""), speaker_note="Mulai dari pain point dan biaya masalah.")
        page += 1
        seed_content_slide(prs, data, page, "Solution", "A simpler way to solve the problem", "Jelaskan perubahan sebelum dan sesudah produk dipakai.", lines(data.get("solution", ""), 5), "Solusi harus langsung menjawab problem utama.", side_title="Value proposition", side_body=data.get("value_prop", ""), speaker_note="Jelaskan before-after customer.")
        page += 1
        seed_content_slide(prs, data, page, "Product", "The product turns workflow into measurable outcomes", "Demo flow utama, bukan semua fitur.", [data.get("product_flow", ""), data.get("product_benefit", "")] + lines(data.get("features", ""), 4), "Investor harus paham cara produk menciptakan nilai.", speaker_note="Jelaskan input -> proses -> output -> dampak.")
        page += 1
        seed_content_slide(prs, data, page, "Market", "A large market with a focused entry wedge", "TAM menunjukkan potensi; SOM menunjukkan fokus eksekusi awal.", lines(data.get("market_notes", ""), 5), "Pasar awal harus spesifik dan bisa dimenangkan.", metric_items=[("TAM", data.get("tam", "")), ("SAM", data.get("sam", "")), ("SOM", data.get("som", ""))])
        page += 1
        seed_content_slide(prs, data, page, "Business Model", "Revenue engine matches the selected model", business_model_template(data).get("description", ""), lines(data.get("business_model", ""), 5) + [business_model_formula_text(data)], "Gunakan metric yang sesuai dengan model bisnis.", metric_items=get_business_model_slide_metrics(data))
        page += 1
        seed_content_slide(prs, data, page, "Traction", "Early demand is already visible", "Tampilkan demand signal paling kuat.", lines(data.get("traction_notes", ""), 5), "Traction harus membuktikan demand, bukan vanity metrics.", metric_items=[("Users", data.get("users", "")), ("Revenue", data.get("revenue", "")), ("Growth", data.get("growth", "")), ("Retention", data.get("retention", ""))])
        page += 1
        seed_content_slide(prs, data, page, "GTM + Competition", "Repeatable acquisition with a clear reason to win", "GTM dan kompetisi dipadatkan untuk durasi 8-10 menit.", [f"ICP: {data.get('icp', '')}", f"Channel: {data.get('channel', '')}"] + lines(data.get("gtm", ""), 3), "Investor perlu melihat channel dan positioning.", side_title="Why we win", side_body="\n".join(compact_competition_lines(data, 3) + [data.get("competition_summary", "")]))
        page += 1
        seed_content_slide(prs, data, page, "Financials + Milestones", "Funding converts into measurable execution", "Proyeksi dan milestones digabung agar alur tetap ringkas.", [f"Revenue Y1/Y2/Y3: {money(data.get('rev1', 0), data.get('currency', 'Rp'))} / {money(data.get('rev2', 0), data.get('currency', 'Rp'))} / {money(data.get('rev3', 0), data.get('currency', 'Rp'))}", f"Runway: {data.get('runway', '')} bulan"] + compact_milestone_lines(data, 3), "Funding harus terhubung langsung dengan milestone.", side_title="Next milestone", side_body=milestone_headline(data))
        page += 1
        seed_content_slide(prs, data, page, "Team", "The right team for this market", "Founder-market fit menjadi sinyal kuat di seed stage.", lines(data.get("team", ""), 5), "Tim harus terlihat punya unfair advantage.", side_title="Founder-market fit", side_body=data.get("founder_fit", ""))
        page += 1
        seed_content_slide(prs, data, page, "Fundraising Ask", "We are raising to reach the next inflection point", "Ask, runway, use of funds, dan next round logic.", lines(data.get("use_of_funds", ""), 5), "Ask harus jelas dan measurable.", side_title=data.get("round", "Round"), side_body=f'{money(data.get("ask", 0), data.get("currency", "Rp"))}\n\nNext milestone: {milestone_headline(data)}\n\n{data.get("next_round", "")}')
        page += 1
        add_compact_closing(prs, data, page)
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output


def build_deck(data, image_buffer=None):
    profile = pitch_duration_profile(data)
    if profile["mode"] == "full":
        return build_full_deck(data, image_buffer)
    return build_compact_deck(data, image_buffer)


# ==============================
# Streamlit UI
# ==============================
st.title("Seed Investor Pitch Deck Generator")
st.caption(
    "Generator PPTX pitch deck profesional untuk seed-stage startup: story-led, data-first, "
    "minim teks, kompetitor dinamis, milestone execution plan, dan analisa investor readiness otomatis."
)

st.markdown(f'<div class="developer-footer">{DEVELOPER_FOOTER}</div>', unsafe_allow_html=True)

with st.sidebar:
    st.header("Brand & Format")
    guide(
        "Pengaturan visual deck",
        "Pilih warna aksen yang sesuai brand. Tampilan aplikasi sudah theme-aware: ketika Streamlit memakai light, dark, atau system theme, warna teks, kartu, input, tab, dan panel analisa akan menyesuaikan agar tetap terbaca.",
    )
    accent_color = st.color_picker(
        "Warna aksen",
        "#2563EB",
        help="Warna ini dipakai untuk headline kecil, metric utama, tabel, dan highlight investor takeaway.",
    )
    currency = st.selectbox(
        "Mata uang",
        ["Rp", "USD"],
        help="Pilih mata uang yang dipakai pada ask funding dan proyeksi finansial.",
    )
    pitch_duration_minutes = st.selectbox(
        "Durasi pitching",
        PITCH_DURATION_OPTIONS,
        index=PITCH_DURATION_OPTIONS.index(10),
        format_func=lambda value: f"{value} menit - {PITCH_DURATION_PROFILES[value]['name']}",
        help="Generator akan menyesuaikan jumlah slide, penggabungan topik, timing PDF, dan kedalaman narasi berdasarkan durasi ini.",
    )
    active_profile = pitch_duration_profile({"pitch_duration_minutes": pitch_duration_minutes})
    st.info(f"Skema: {active_profile['scheme']}")
    include_insight_slide = st.checkbox(
        "Tambahkan slide analisa otomatis",
        value=True,
        help="Jika aktif, PPTX akan menyertakan slide Investor Readiness berisi insight dari data yang diinput.",
    )
    st.success(DEVELOPER_FOOTER)

identity, story, market, finance, team_asset, glossary_tab, insight_tab = st.tabs(
    [
        "Identitas",
        "Story",
        "Market, Model & Traction",
        "Financial & Funding",
        "Team & Competition",
        "Istilah & Rumus",
        "Analisa",
    ]
)

with identity:
    guide(
        "Identitas deck",
        "Bagian ini membentuk slide pembuka. Investor harus langsung paham: nama startup, siapa yang presentasi, startup menyelesaikan masalah apa, dan berapa pendanaan yang sedang dicari.",
    )

    col1, col2 = st.columns(2)

    with col1:
        company = st.text_input(
            "Nama startup",
            "Rex AI",
            help="Tulis nama legal/brand startup. Nama ini juga muncul di footer setiap slide.",
        )
        one_liner = st.text_input(
            "One-liner",
            "AI copilot untuk membantu UMKM membuat laporan keuangan otomatis.",
            help="Formula praktis: [produk] untuk [customer] agar [hasil bisnis]. Hindari slogan abstrak.",
        )
        presenter = st.text_input(
            "Presenter",
            "Rex Founder",
            help="Nama founder atau orang yang akan melakukan pitching.",
        )

    with col2:
        contact = st.text_input(
            "Kontak",
            "founder@rex.ai",
            help="Email atau kontak bisnis yang akan muncul di cover dan closing slide.",
        )
        round_name = st.text_input(
            "Round",
            "Seed Round",
            help="Contoh: Pre-Seed Round, Seed Round, Bridge Round. Sesuaikan dengan tahap fundraising.",
        )
        ask = st.number_input(
            "Jumlah pendanaan",
            min_value=0,
            value=1_500_000_000,
            step=50_000_000,
            help="Jumlah dana yang dicari. Pastikan selaras dengan runway dan milestone 12-18 bulan.",
        )

with story:
    guide(
        "Story: problem, solution, dan product",
        "Bagian ini adalah jantung pitch. Investor perlu melihat masalah yang urgent, solusi yang tajam, dan produk yang bisa menjawab masalah dengan cara lebih baik dari alternatif yang ada.",
    )

    col1, col2 = st.columns(2)

    with col1:
        problem = st.text_area(
            "Problem",
            "UMKM sulit mencatat transaksi harian.\nLaporan keuangan masih manual dan rawan salah.\nPemilik usaha tidak punya insight real-time.",
            height=130,
            help="Tulis 3-5 pain point utama. Fokus pada masalah yang sering, mahal, dan penting untuk customer.",
        )
        problem_evidence = st.text_area(
            "Evidence problem",
            "Interview awal: 24 dari 30 UMKM masih memakai catatan manual atau spreadsheet.",
            height=90,
            help="Bukti bahwa problem nyata: interview, pilot, data market, waiting list, churn dari solusi lama, atau biaya masalah.",
        )
        solution = st.text_area(
            "Solution",
            "Input transaksi via chat.\nOCR struk otomatis.\nDashboard cashflow dan laba rugi.\nRekomendasi aksi berbasis data.",
            height=130,
            help="Jelaskan cara solusi bekerja dalam 3-5 poin. Hindari terlalu banyak fitur teknis.",
        )
        value_prop = st.text_area(
            "Value proposition",
            "Menghemat waktu pencatatan dan membantu owner memahami kondisi keuangan tanpa akuntan penuh waktu.",
            height=90,
            help="Tulis manfaat utama yang dirasakan customer: hemat waktu, naik revenue, turun biaya, lebih cepat, lebih akurat, atau lebih mudah.",
        )

    with col2:
        product_flow = st.text_area(
            "Cara kerja produk",
            "User kirim transaksi/foto struk → sistem klasifikasi otomatis → dashboard dan laporan dibuat real-time.",
            height=90,
            help="Tulis alur produk dari input → proses → output → dampak. Ini membantu saat demo pitching.",
        )
        product_benefit = st.text_area(
            "Benefit produk",
            "Owner melihat omzet, margin, utang-piutang, dan cashflow dalam satu tempat.",
            height=90,
            help="Tulis outcome produk, bukan fitur. Contoh: keputusan lebih cepat, error turun, cashflow lebih jelas.",
        )
        features = st.text_area(
            "Fitur utama",
            "OCR struk\nAuto-categorization\nCashflow dashboard\nExport laporan PDF/Excel\nReminder piutang",
            height=130,
            help="Pilih maksimal 5 fitur yang paling mendukung value proposition. Jangan jadikan slide sebagai daftar fitur panjang.",
        )
        product_image = st.file_uploader(
            "Screenshot/mockup produk",
            type=["png", "jpg", "jpeg"],
            help="Upload mockup, screenshot, atau visual produk. Untuk seed deck, visual produk sangat membantu investor memahami demo secara cepat.",
        )

with market:
    guide(
        "Market, business model, traction, dan GTM",
        "Bagian ini menjawab: seberapa besar peluangnya, siapa customer awal, bagaimana startup menghasilkan uang, apakah ada demand awal, dan bagaimana akuisisi pelanggan akan diulang.",
    )

    col1, col2 = st.columns(2)

    with col1:
        tam = st.text_input("TAM", "Rp 120T", help="Total Addressable Market: keseluruhan peluang pasar jika startup menang besar.")
        sam = st.text_input("SAM", "Rp 18T", help="Serviceable Available Market: pasar yang realistis dijangkau oleh produk dan model bisnis saat ini.")
        som = st.text_input("SOM", "Rp 450M", help="Serviceable Obtainable Market: target pasar awal yang realistis dimenangkan dalam beberapa tahun pertama.")
        market_notes = st.text_area(
            "Market notes",
            "Target awal: UMKM F&B dan ritel.\nWedge: bisnis yang sudah memakai WhatsApp untuk operasional.\nEkspansi: inventory, payroll, dan embedded financing.",
            height=120,
            help="Jelaskan segmen awal, alasan memilih segmen tersebut, dan rencana ekspansi pasar.",
        )
        business_model_type = st.selectbox(
            "Jenis model bisnis",
            list(BUSINESS_MODEL_TEMPLATES.keys()),
            index=0,
            help="Pilih model bisnis utama. Aplikasi akan menyesuaikan metrik, format slide Business Model, insight, dan PDF skenario.",
        )
        selected_model_template = business_model_template(business_model_type)
        st.markdown(
            f"""
            <div class="readable-panel">
                <p><strong>Format model:</strong> {html.escape(selected_model_template['description'])}</p>
                <p><strong>Cocok untuk:</strong> {html.escape(selected_model_template['best_for'])}</p>
                <p><strong>Yang harus ditekankan saat pitching:</strong> {html.escape(selected_model_template['pitch_focus'])}</p>
                <p><strong>Rumus praktis:</strong> {html.escape(selected_model_template['formula'])}</p>
            </div>
            """,
            unsafe_allow_html=True,
        )

        # Metrik spesifik model bisnis.
        # Variabel ini wajib dibuat sebelum dictionary `data`, karena dipakai oleh
        # slide Business Model, analisa, dan PDF scenario guide.
        default_model_metrics = selected_model_template.get("metrics", [])[:4]
        model_metric_labels = []
        model_metric_values = []

        with st.expander("Metrik utama sesuai model bisnis", expanded=True):
            st.caption(
                "Sesuaikan label dan nilai metrik agar cocok dengan revenue engine startup. "
                "Contoh: SaaS memakai MRR/ARR, marketplace memakai GMV & take rate, "
                "e-commerce memakai AOV & repeat purchase."
            )

            for metric_idx, metric in enumerate(default_model_metrics):
                default_label = metric[0] if len(metric) > 0 else f"Metric {metric_idx + 1}"
                default_value = metric[1] if len(metric) > 1 else ""

                metric_col_1, metric_col_2 = st.columns([1, 1])

                with metric_col_1:
                    metric_label = st.text_input(
                        f"Label metrik {metric_idx + 1}",
                        value=default_label,
                        key=f"model_metric_label_{metric_idx}",
                        help="Nama metrik yang akan muncul di slide Business Model, misalnya MRR, GMV, Take Rate, AOV, NRR, atau Transaction Volume.",
                    )

                with metric_col_2:
                    metric_value = st.text_input(
                        f"Nilai metrik {metric_idx + 1}",
                        value=default_value,
                        key=f"model_metric_value_{metric_idx}",
                        help="Nilai metrik yang akan ditampilkan di deck. Gunakan format yang mudah dibaca investor, misalnya Rp 85 juta MRR, 5% take rate, atau 2 juta API calls/bulan.",
                    )

                model_metric_labels.append(metric_label)
                model_metric_values.append(metric_value)

        arpu = st.text_input("ARPU", "Rp 99.000/bulan", help="Average Revenue Per User. Tulis pendapatan rata-rata per customer/user.")
        gross_margin = st.text_input("Gross margin", "78%", help="Margin kotor. Untuk software/SaaS, margin tinggi menjadi sinyal scalability.")
        cac = st.text_input("CAC", "Rp 140.000", help="Customer Acquisition Cost. Tulis biaya rata-rata untuk mendapatkan satu customer.")
        payback = st.text_input("Payback", "< 2 bulan", help="Periode balik modal CAC. Semakin pendek, semakin kuat unit economics.")
        business_model = st.text_area(
            "Business model",
            "Subscription SaaS bulanan.\nPaket Pro untuk multi-outlet.\nAdd-on laporan pajak dan inventory.\nRevenue share dari pembiayaan UMKM.",
            height=120,
            help="Jelaskan siapa yang membayar, berapa, kapan membayar, dan potensi expansion revenue.",
        )

    with col2:
        users = st.text_input("Users/customers", "1.200 users", help="Jumlah user/customer aktif, customer berbayar, pilot, atau akun terdaftar. Sebutkan konteksnya.")
        revenue = st.text_input("Revenue/GMV", "Rp 85 juta MRR", help="Revenue bulanan, ARR, GMV, atau pipeline. Untuk investor, revenue berulang lebih kuat daripada vanity metrics.")
        growth = st.text_input("Growth", "+28% MoM", help="Pertumbuhan month-over-month atau quarter-over-quarter. Jelaskan basis perhitungannya.")
        retention = st.text_input("Retention", "72% D30", help="Retention D7/D30, logo retention, revenue retention, atau repeat usage. Retention menjawab apakah produk benar-benar dipakai ulang.")
        traction_notes = st.text_area(
            "Traction notes",
            "Pilot berbayar dengan 35 UMKM.\nPipeline 180 UMKM dari partner.\nChurn turun setelah fitur reminder piutang dirilis.",
            height=120,
            help="Tulis bukti demand: revenue, active usage, LOI, pilot, partnership, pipeline, atau waiting list.",
        )
        icp = st.text_area(
            "ICP",
            "Pemilik UMKM F&B/ritel dengan 3-50 transaksi per hari.",
            height=80,
            help="Ideal Customer Profile. Spesifikkan jenis customer yang paling sakit masalahnya dan paling mudah dijangkau.",
        )
        channel = st.text_area(
            "Channel utama",
            "Komunitas UMKM, reseller akuntansi, konten edukasi, partnership POS.",
            height=80,
            help="Channel akuisisi utama: komunitas, sales, referral, partnership, content, outbound, marketplace, atau channel lain.",
        )
        gtm = st.text_area(
            "GTM",
            "Akuisisi lewat komunitas UMKM.\nReferral program.\nPartnership konsultan pajak dan POS.\nKonten edukasi cashflow sebagai lead magnet.",
            height=110,
            help="Go-To-Market: jelaskan proses akuisisi dari lead source, conversion, onboarding, sampai retention.",
        )

with finance:
    guide(
        "Financial, funding ask, dan milestone",
        "Bagian ini membantu investor melihat apakah jumlah pendanaan masuk akal. Hubungkan funding dengan runway, penggunaan dana, dan milestone yang bisa dicapai dalam 12-18 bulan.",
    )

    col1, col2 = st.columns(2)

    with col1:
        rev1 = st.number_input("Revenue Y1", value=1_200_000_000, step=100_000_000, help="Proyeksi revenue tahun pertama setelah funding. Gunakan asumsi realistis.")
        rev2 = st.number_input("Revenue Y2", value=4_800_000_000, step=100_000_000, help="Proyeksi revenue tahun kedua. Pastikan growth didukung asumsi customer, ARPU, dan channel.")
        rev3 = st.number_input("Revenue Y3", value=13_500_000_000, step=100_000_000, help="Proyeksi revenue tahun ketiga. Gunakan untuk menunjukkan potensi skala.")
        cost1 = st.number_input("Operating Cost Y1", value=2_000_000_000, step=100_000_000, help="Biaya operasional tahun pertama: tim, produk, marketing, sales, tools, legal, operasional.")
        cost2 = st.number_input("Operating Cost Y2", value=5_200_000_000, step=100_000_000, help="Biaya operasional tahun kedua.")
        cost3 = st.number_input("Operating Cost Y3", value=10_000_000_000, step=100_000_000, help="Biaya operasional tahun ketiga.")

    with col2:
        profit1 = st.number_input("EBITDA/Profit Y1", value=-800_000_000, step=100_000_000, help="EBITDA/profit tahun pertama. Untuk seed, negatif masih wajar jika diarahkan ke growth.")
        profit2 = st.number_input("EBITDA/Profit Y2", value=-400_000_000, step=100_000_000, help="EBITDA/profit tahun kedua.")
        profit3 = st.number_input("EBITDA/Profit Y3", value=3_500_000_000, step=100_000_000, help="EBITDA/profit tahun ketiga. Ini membantu menunjukkan operating leverage.")
        runway = st.number_input("Runway / bulan", min_value=1, value=18, help="Berapa bulan startup bisa berjalan dengan dana yang dicari. Seed biasanya ditargetkan 12-18 bulan, bergantung strategi.")
        milestone = st.text_input(
            "Milestone utama / headline",
            "Rp 500 juta MRR, 8.000 active businesses, gross margin >80%",
            help="Ringkasan milestone paling penting yang muncul di slide Financials dan Ask.",
        )
        use_of_funds = st.text_area(
            "Use of funds",
            "Product & engineering: 40%\nSales & marketing: 35%\nCustomer success: 15%\nOperations & legal: 10%",
            height=110,
            help="Pembagian penggunaan dana. Buat logis dan langsung terhubung ke milestone.",
        )
        next_round = st.text_area(
            "Next round logic",
            "Raise Seed/Series A setelah akuisisi repeatable dan retention terbukti.",
            height=80,
            help="Jelaskan kondisi apa yang membuat startup siap raise round berikutnya.",
        )

    guide(
        "Detail milestone",
        "Tambahkan milestone bertahap agar investor melihat rute eksekusi. Setiap milestone sebaiknya punya periode, target, success metric, dan owner. Jika lebih dari 4 milestone, PPT akan otomatis memecahnya menjadi beberapa slide.",
    )

    milestone_defaults = [
        {
            "period": "0-3 bulan",
            "target": "Rilis MVP stabil dan onboarding 100 customer awal",
            "metric": "100 active customers, activation >60%",
            "owner": "Product & Growth",
        },
        {
            "period": "4-6 bulan",
            "target": "Validasi channel akuisisi utama",
            "metric": "CAC payback <3 bulan, 500 customers",
            "owner": "Growth",
        },
        {
            "period": "7-12 bulan",
            "target": "Capai repeatable revenue motion",
            "metric": "Rp 250 juta MRR, retention D30 >70%",
            "owner": "Sales & CS",
        },
        {
            "period": "13-18 bulan",
            "target": "Siap raise next round dengan traction kuat",
            "metric": "Rp 500 juta MRR, gross margin >80%",
            "owner": "Leadership",
        },
    ]

    milestone_count = int(
        st.number_input(
            "Jumlah milestone detail",
            min_value=1,
            max_value=8,
            value=4,
            step=1,
            help="PPT akan menyesuaikan layout otomatis. Maksimal 4 milestone per slide.",
        )
    )

    milestones = []

    for idx in range(milestone_count):
        default = milestone_defaults[idx] if idx < len(milestone_defaults) else {
            "period": f"Tahap {idx + 1}",
            "target": "Target milestone",
            "metric": "Metric keberhasilan",
            "owner": "Team",
        }

        with st.expander(f"Milestone {idx + 1}: {default['period']}", expanded=idx < 2):
            m1, m2 = st.columns(2)
            with m1:
                period = st.text_input(
                    f"Periode milestone {idx + 1}",
                    default["period"],
                    key=f"milestone_period_{idx}",
                    help="Contoh: 0-3 bulan, Q1 2027, sebelum Seed extension, sebelum Series A.",
                )
                target = st.text_area(
                    f"Target milestone {idx + 1}",
                    default["target"],
                    height=80,
                    key=f"milestone_target_{idx}",
                    help="Target eksekusi yang ingin dicapai pada periode ini.",
                )
            with m2:
                metric = st.text_input(
                    f"Success metric {idx + 1}",
                    default["metric"],
                    key=f"milestone_metric_{idx}",
                    help="Buat measurable: MRR, active user, retention, CAC payback, gross margin, jumlah partnership, atau product release.",
                )
                owner = st.text_input(
                    f"Owner milestone {idx + 1}",
                    default["owner"],
                    key=f"milestone_owner_{idx}",
                    help="Tim/role yang bertanggung jawab: Product, Growth, Sales, CS, Engineering, atau Leadership.",
                )

        milestones.append(
            {
                "period": period,
                "target": target,
                "metric": metric,
                "owner": owner,
            }
        )

with team_asset:
    guide(
        "Team dan competition",
        "Investor seed sangat memperhatikan founder-market fit dan alasan startup ini bisa menang. Tambahkan kompetitor langsung, tidak langsung, dan status quo. PPT akan otomatis menyesuaikan layout jika kompetitor lebih dari satu atau lebih dari empat.",
    )

    col1, col2 = st.columns(2)

    with col1:
        team = st.text_area(
            "Team",
            "Rex - CEO - 6 tahun membangun SaaS UMKM.\nNadia - CTO - ex fintech data engineer.\nBima - Growth - pernah scale komunitas UMKM 50k members.",
            height=120,
            help="Tulis nama, role, dan pengalaman relevan. Fokus pada kemampuan yang langsung berhubungan dengan problem, produk, dan market.",
        )
        founder_fit = st.text_area(
            "Founder-market fit",
            "Tim pernah membangun tools operasional untuk UMKM dan punya akses langsung ke komunitas target.",
            height=90,
            help="Jelaskan unfair advantage tim: pengalaman domain, akses distribusi, technical edge, atau insight unik dari market.",
        )
        closing = st.text_input(
            "Closing line",
            "Let’s help millions of small businesses understand their money.",
            help="Kalimat penutup yang merangkum visi besar startup.",
        )

    with col2:
        competition_summary = st.text_area(
            "Ringkasan kompetisi / narrative advantage",
            "Kami menang di simplicity, distribusi lokal, dan data workflow UMKM.",
            height=100,
            help="Satu kalimat narrative advantage: kenapa startup ini bisa menang dibanding alternatif yang sudah dipakai customer.",
        )
        competitor_count = int(
            st.number_input(
                "Jumlah kompetitor / alternatif",
                min_value=1,
                max_value=10,
                value=4,
                step=1,
                help="Masukkan kompetitor langsung, tidak langsung, dan status quo. Maksimal 4 kompetitor per slide; jika lebih, PPT otomatis memecah slide.",
            )
        )

    competitor_defaults = [
        {
            "name": "Aplikasi Akuntansi A",
            "category": "Direct competitor",
            "weakness": "Terlalu kompleks untuk mikro-UMKM",
            "advantage": "Workflow chat-first dan onboarding cepat",
        },
        {
            "name": "Spreadsheet/manual",
            "category": "Indirect competitor",
            "weakness": "Tidak real-time dan rawan error",
            "advantage": "Otomatisasi laporan dan insight",
        },
        {
            "name": "Status quo",
            "category": "Status quo",
            "weakness": "Owner tidak punya data untuk keputusan cepat",
            "advantage": "Data harian langsung menjadi rekomendasi aksi",
        },
        {
            "name": "POS tanpa analitik keuangan",
            "category": "Adjacent tool",
            "weakness": "Mencatat transaksi tetapi tidak memberi insight cashflow",
            "advantage": "Mengubah transaksi menjadi laporan dan rekomendasi bisnis",
        },
    ]

    competitors = []

    st.markdown("### Peta kompetitor")
    st.caption(
        "Untuk deck seed, kompetitor tidak harus hanya produk yang sama. Masukkan juga cara manual, spreadsheet, agency, status quo, atau adjacent tool yang saat ini dipakai customer."
    )

    for idx in range(competitor_count):
        default = competitor_defaults[idx] if idx < len(competitor_defaults) else {
            "name": f"Kompetitor {idx + 1}",
            "category": "Alternative",
            "weakness": "Kelemahan dari sudut pandang customer",
            "advantage": "Keunggulan produk Anda",
        }

        with st.expander(f"Kompetitor / alternatif {idx + 1}: {default['name']}", expanded=idx < 3):
            c1, c2 = st.columns(2)
            with c1:
                comp_name = st.text_input(
                    f"Nama kompetitor / alternatif {idx + 1}",
                    default["name"],
                    key=f"competitor_name_{idx}",
                    help="Nama produk, perusahaan, workflow manual, spreadsheet, agency, atau status quo.",
                )
                comp_category = st.selectbox(
                    f"Kategori {idx + 1}",
                    ["Direct competitor", "Indirect competitor", "Status quo", "Adjacent tool", "Alternative"],
                    index=["Direct competitor", "Indirect competitor", "Status quo", "Adjacent tool", "Alternative"].index(default["category"]) if default["category"] in ["Direct competitor", "Indirect competitor", "Status quo", "Adjacent tool", "Alternative"] else 4,
                    key=f"competitor_category_{idx}",
                    help="Kategorisasi membantu investor memahami medan kompetisi.",
                )
            with c2:
                comp_weakness = st.text_area(
                    f"Kelemahan {idx + 1}",
                    default["weakness"],
                    height=80,
                    key=f"competitor_weakness_{idx}",
                    help="Tulis kelemahan dari perspektif customer: mahal, lambat, rumit, tidak real-time, sulit diadopsi, tidak terintegrasi, dsb.",
                )
                comp_advantage = st.text_area(
                    f"Keunggulan kita vs alternatif {idx + 1}",
                    default["advantage"],
                    height=80,
                    key=f"competitor_advantage_{idx}",
                    help="Tulis keunggulan spesifik dan defensible: distribusi, data, UX, cost, speed, domain expertise, switching workflow, dsb.",
                )

        competitors.append(
            {
                "name": comp_name,
                "category": comp_category,
                "weakness": comp_weakness,
                "advantage": comp_advantage,
            }
        )

with glossary_tab:
    render_glossary_section()


# Collect current data for analysis and generation.
data = {
    "company": company,
    "one_liner": one_liner,
    "presenter": presenter,
    "contact": contact,
    "round": round_name,
    "ask": ask,
    "currency": currency,
    "pitch_duration_minutes": pitch_duration_minutes,
    "accent_color": accent_color,
    "include_insight_slide": include_insight_slide,
    "problem": problem,
    "problem_evidence": problem_evidence,
    "solution": solution,
    "value_prop": value_prop,
    "product_flow": product_flow,
    "product_benefit": product_benefit,
    "features": features,
    "tam": tam,
    "sam": sam,
    "som": som,
    "market_notes": market_notes,
    "arpu": arpu,
    "gross_margin": gross_margin,
    "cac": cac,
    "payback": payback,
    "business_model_type": business_model_type,
    "business_model": business_model,
    "model_metric_labels": model_metric_labels,
    "model_metric_values": model_metric_values,
    "users": users,
    "revenue": revenue,
    "growth": growth,
    "retention": retention,
    "traction_notes": traction_notes,
    "icp": icp,
    "channel": channel,
    "gtm": gtm,
    "competitors": competitors,
    "competitor_1": competitors[0]["name"] if len(competitors) > 0 else "",
    "weakness_1": competitors[0]["weakness"] if len(competitors) > 0 else "",
    "advantage_1": competitors[0]["advantage"] if len(competitors) > 0 else "",
    "competitor_2": competitors[1]["name"] if len(competitors) > 1 else "",
    "weakness_2": competitors[1]["weakness"] if len(competitors) > 1 else "",
    "advantage_2": competitors[1]["advantage"] if len(competitors) > 1 else "",
    "status_quo": competitors[2]["weakness"] if len(competitors) > 2 else "",
    "advantage_3": competitors[2]["advantage"] if len(competitors) > 2 else "",
    "competition_summary": competition_summary,
    "rev1": rev1,
    "rev2": rev2,
    "rev3": rev3,
    "cost1": cost1,
    "cost2": cost2,
    "cost3": cost3,
    "profit1": profit1,
    "profit2": profit2,
    "profit3": profit3,
    "runway": runway,
    "milestone": milestone,
    "milestones": milestones,
    "use_of_funds": use_of_funds,
    "next_round": next_round,
    "team": team,
    "founder_fit": founder_fit,
    "closing": closing,
}

with insight_tab:
    guide(
        "Analisa otomatis dari data pitching",
        "Bagian ini tidak perlu diisi manual. Sistem membaca input Anda dan memberi ringkasan kekuatan, risiko pertanyaan investor, serta rekomendasi perbaikan narasi sebelum deck di-download.",
    )
    show_insights(generate_investor_insights(data))

st.divider()

col_generate, col_hint = st.columns([1, 2])

with col_generate:
    generate = st.button("Generate Seed Investor Pitch Deck", type="primary", use_container_width=True)

with col_hint:
    st.caption(
        "Sebelum generate, cek tab Analisa untuk memastikan narasi problem, traction, financial, dan ask sudah konsisten. "
        "Gunakan tab Istilah & Rumus untuk memahami cara menghitung metrik sebelum mengisi angka. "
        "PPTX dan PDF scenario guide akan menyesuaikan durasi pitching, model bisnis, dan otomatis membawa footer Developed by Galuh Adi Insani."
    )

if generate:
    if not company.strip():
        st.error("Nama startup wajib diisi.")
        st.stop()

    image_buffer = BytesIO(product_image.read()) if product_image else None
    pptx = build_deck(data, image_buffer)
    scenario_pdf = build_scenario_pdf(data)

    st.success("Pitch deck dan PDF scenario guide berhasil dibuat.")
    show_insights(generate_investor_insights(data))

    st.markdown(
        """
        <div class="readable-panel">
            <p><strong>PDF Pitch Scenario Guide</strong></p>
            <p>PDF ini mengikuti urutan slide PPTX dan berisi tujuan slide, narasi bicara, transisi, pertanyaan investor, kamus istilah startup, rumus perhitungan, contoh angka, serta checklist latihan pitching.</p>
        </div>
        """,
        unsafe_allow_html=True,
    )

    download_col1, download_col2 = st.columns(2)

    with download_col1:
        st.download_button(
            "📥 Download Seed Investor Pitch Deck (.pptx)",
            data=pptx,
            file_name=f"{filename(company)}-seed-investor-pitch-deck.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True,
        )

    with download_col2:
        st.download_button(
            "📘 Download Pitch Scenario Guide (.pdf)",
            data=scenario_pdf,
            file_name=f"{filename(company)}-pitch-scenario-guide.pdf",
            mime="application/pdf",
            use_container_width=True,
        )
